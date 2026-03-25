#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation for Antah Asti Prarambh project.
Comparative Structural Proteomics of Chaperonin Substrates.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === CONFIGURATION ===
BASE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE, "results/phase2/figures")
OUT_FILE = os.path.join(BASE, "Antah_Asti_Prarambh_Presentation.pptx")

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x86, 0xC8)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_GREEN = RGBColor(0x2D, 0x8E, 0x4E)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
VERY_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === HELPER FUNCTIONS ===

def add_background(slide, color=WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=BLACK, spacing=Pt(6), bold_prefix=False):
    """Add bulleted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = prefix + ": "
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox

def add_header_bar(slide, title_text, subtitle_text=None):
    """Add a styled header bar at the top of a slide."""
    bar = add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.55), title_text,
                font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.62), Inches(12), Inches(0.4), subtitle_text,
                    font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD))

def add_footer(slide, slide_num, total=30):
    """Add footer with slide number."""
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
                f"{slide_num}/{total}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(5), Inches(0.3),
                "Antah Asti Prarambh | Vishal Bharti, CSIR-IGIB", font_size=10, color=GRAY)

def add_table(slide, left, top, width, height, data, col_widths=None, header_color=DARK_BLUE):
    """Add a styled table."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = BLACK

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERY_LIGHT_BLUE

    return table_shape

# ======================================================================
# SLIDE 1: TITLE SLIDE
# ======================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, DARK_BLUE)

# Accent bar
add_shape(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(1), Inches(1.0), Inches(11.333), Inches(0.7),
            "ANTAH ASTI PRARAMBH", font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(0.5),
            '"The End is the Beginning"', font_size=22, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.8),
            "A Comparative Structural Proteomics Study of\nGroup I Chaperonin Substrates",
            font_size=26, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(0.5),
            "GroEL (E. coli) vs HSP60 (Human Mitochondria)",
            font_size=18, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.4),
            "Vishal Bharti", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.95), Inches(11.333), Inches(0.4),
            "CSIR-Institute of Genomics and Integrative Biology (IGIB), New Delhi",
            font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.4), Inches(11.333), Inches(0.4),
            "March 2026", font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

sn = 1

# ======================================================================
# SLIDE 2: OUTLINE
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Presentation Outline")
add_footer(slide, sn)

items_left = [
    "1.  Background & Motivation",
    "2.  The Three Scientific Goals",
    "3.  The Seven Datasets",
    "4.  Critical Methodological Decisions",
    "5.  Computational Workflow",
    "6.  Phase 1: Pilot Analysis (1,390 proteins)",
    "7.  Phase 2: Full-Scale HPC (25,007 proteins)",
]
items_right = [
    "8.   Results: Domain Architecture",
    "9.   Results: N-vs-C Structural Asymmetry",
    "10. Results: Mitochondrial Targeting",
    "11. Results: Cross-Species Conservation",
    "12. Biological Synthesis",
    "13. Limitations & Future Work",
    "14. Conclusions",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), items_left, font_size=17, spacing=Pt(10))
add_bullet_slide(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.5), items_right, font_size=17, spacing=Pt(10))

# ======================================================================
# SLIDE 3: BACKGROUND & MOTIVATION
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Background: Chaperonins in Protein Folding", "Why do some proteins need help to fold?")
add_footer(slide, sn)

# Left column: GroEL
add_shape(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.3), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.4), Inches(0.4),
            "GroEL/GroES (E. coli)", font_size=20, bold=True, color=MED_BLUE)
items = [
    "~800 kDa double-ring complex",
    "Assists ~250 substrate proteins (10-15% of proteome)",
    "Kerner et al. 2005 — gold standard classification:",
    "    Class I (38): Spontaneous folders, accelerated by GroEL",
    "    Class II (126): Partially dependent — fold slowly without GroEL",
    "    Class III (84): Obligate — cannot fold without GroEL",
    "Co-chaperonin: GroES (lid)",
]
add_bullet_slide(slide, Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.5), items, font_size=14, spacing=Pt(5))

# Right column: HSP60
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.3), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.4),
            "HSP60/HSP10 (Human Mitochondria)", font_size=20, bold=True, color=ACCENT_ORANGE)
items = [
    "Mitochondrial matrix chaperonin",
    "266 Tier-1 substrates (Morten et al. 2020, SILAC co-IP)",
    "Nuclear-encoded proteins imported into matrix",
    "Must unfold for import through TOM/TIM complexes",
    "Re-fold in matrix with HSP60 assistance",
    "Co-chaperonin: HSP10 (HSPE1)",
    "Evolutionary homolog of GroEL (~2 billion years divergence)",
]
add_bullet_slide(slide, Inches(7.2), Inches(2.0), Inches(5.4), Inches(4.5), items, font_size=14, spacing=Pt(5))

# Central arrow
add_textbox(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(1.0),
            "~2 billion years\nof evolution\n(endosymbiosis)", font_size=12, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 4: WHY COMPARE / PROJECT NAME
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Why Compare GroEL and HSP60?", "A natural experiment spanning 2 billion years")
add_footer(slide, sn)

items = [
    "If substrate properties are CONSERVED:",
    "    Strong evolutionary constraints on what makes a protein need chaperonin help",
    "    Structural features making folding difficult are fundamental & ancient",
    "",
    "If substrate properties have DIVERGED:",
    "    Chaperonin-substrate co-evolution during transition from bacterium to organelle",
    "    Mitochondrial import adds new constraints (MTS, unfolding, re-folding)",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.4), Inches(7), Inches(3.5), items, font_size=16, spacing=Pt(6))

# Name explanation box
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.4),
            'Project Name: "Antah Asti Prarambh" (Sanskrit: "The End is the Beginning")',
            font_size=18, bold=True, color=ACCENT_ORANGE)
items = [
    "Central question: Does the C-terminus ('end' of translation) represent a new 'beginning' for chaperonin-assisted folding?",
    "Proteins are synthesized N-to-C on ribosomes. N-terminal regions emerge & fold first.",
    "C-terminal regions emerge last — potentially the region most needing chaperonin help.",
    "We test whether this N-vs-C asymmetry is specific to chaperonin substrates.",
]
add_bullet_slide(slide, Inches(0.8), Inches(5.35), Inches(11.5), Inches(1.5), items, font_size=13, spacing=Pt(3))

# ======================================================================
# SLIDE 5: THREE SCIENTIFIC GOALS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Three Scientific Goals")
add_footer(slide, sn)

# Goal 1
add_shape(slide, Inches(0.4), Inches(1.4), Inches(4.0), Inches(2.5), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(3.6), Inches(0.4),
            "Goal 1: Domain Architecture", font_size=17, bold=True, color=MED_BLUE)
items = [
    "Do chaperonin substrates have distinctive structural folds?",
    "Compare CATH superfamily distributions",
    "Substrates vs size-matched proteome background",
    "Fisher's exact test with odds ratios",
]
add_bullet_slide(slide, Inches(0.6), Inches(2.0), Inches(3.6), Inches(1.8), items, font_size=12, spacing=Pt(3))

# Goal 2
add_shape(slide, Inches(4.7), Inches(1.4), Inches(4.0), Inches(2.5), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(4.9), Inches(1.5), Inches(3.6), Inches(0.4),
            "Goal 2: N-vs-C Asymmetry", font_size=17, bold=True, color=ACCENT_ORANGE)
items = [
    "Are N-terminal domains more complex (higher contact order)?",
    "Three-region decomposition: pre-domain, N-domain, C-region",
    "Paired within-protein comparisons",
    "Substrate-specific or universal?",
]
add_bullet_slide(slide, Inches(4.9), Inches(2.0), Inches(3.6), Inches(1.8), items, font_size=12, spacing=Pt(3))

# Goal 3
add_shape(slide, Inches(9.0), Inches(1.4), Inches(4.0), Inches(2.5), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(9.2), Inches(1.5), Inches(3.6), Inches(0.4),
            "Goal 3: MTS Architecture", font_size=17, bold=True, color=ACCENT_GREEN)
items = [
    "How do targeting signals relate to structural domains?",
    "HSP60 substrates: matrix enrichment?",
    "MTS = pre-domain extension or domain overlap?",
    "The 'landing pad' model",
]
add_bullet_slide(slide, Inches(9.2), Inches(2.0), Inches(3.6), Inches(1.8), items, font_size=12, spacing=Pt(3))

# Pre-registered hypotheses section
add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
            "Pre-Registered Hypotheses (9 total, 3 families)", font_size=18, bold=True, color=DARK_BLUE)

hyp_data = [
    ["Family", "ID", "Hypothesis", "Test", "Effect Size"],
    ["Domain\nArchitecture", "H1.1", "GroEL substrates enriched for specific superfamilies (TIM barrels)", "Fisher's exact", "Odds ratio"],
    ["", "H1.2", "HSP60 substrates show fold enrichment vs matrix background", "Fisher's exact", "Odds ratio"],
    ["", "H1.3", "Fold enrichment conserved between GroEL and HSP60", "Chi-squared", "Cramer's V"],
    ["N-vs-C\nAsymmetry", "H2.1", "N-domains have different contact order than C-regions", "Wilcoxon signed-rank", "Rank-biserial r"],
    ["", "H2.2", "N-vs-C asymmetry GREATER in substrates vs background", "Mann-Whitney U", "Rank-biserial r"],
    ["", "H2.3", "Class III > Class I asymmetry (GroEL class gradient)", "Kruskal-Wallis H", "Eta-squared"],
    ["MTS\nTargeting", "H3.1", "HSP60 substrates enriched for matrix localization", "Fisher's exact", "Odds ratio"],
    ["", "H3.2", "MTS-bearing substrates have distinct first-domain properties", "Mann-Whitney U", "Rank-biserial r"],
    ["", "H3.3", "MTS is predominantly a pre-domain extension", "Binomial test", "Proportion"],
]
add_table(slide, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.7), hyp_data,
          col_widths=[Inches(1.3), Inches(0.6), Inches(5.5), Inches(2.5), Inches(2.6)])

# ======================================================================
# SLIDE 6: SEVEN DATASETS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "The Seven Datasets", "Carefully curated datasets enabling all three scientific goals")
add_footer(slide, sn)

ds_data = [
    ["#", "Dataset", "Size", "Source", "Purpose"],
    ["1", "E. coli K-12 Proteome", "4,403 proteins", "UniProt UP000000625", "Background for GroEL comparisons"],
    ["2", "Human Proteome", "20,416 proteins", "UniProt UP000005640", "Parent set; orthology analysis"],
    ["3", "Human Mito Proteome", "1,136 proteins", "MitoCarta 3.0", "Compartment-matched background for HSP60"],
    ["4", "GroEL Substrates", "252 proteins", "Kerner et al. 2005 (Table S3)", "Primary GroEL substrate set (3 classes)"],
    ["5", "HSP60 Tier-1 Substrates", "266 proteins", "Morten et al. 2020 (SILAC co-IP)", "High-confidence HSP60 substrates"],
    ["6", "Cross-Species Homologs", "69 pairs", "OrthoFinder + RBH (computed)", "Conservation analysis across species"],
    ["7", "Mito Matrix Subset", "525 proteins", "MitoCarta 3.0 (matrix only)", "Tightest background for HSP60"],
]
add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.0), ds_data,
          col_widths=[Inches(0.4), Inches(2.5), Inches(1.5), Inches(3.5), Inches(4.8)])

# Data cleaning notes
add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
            "Key Data Cleaning Steps", font_size=16, bold=True, color=DARK_BLUE)
items = [
    "GroEL: 149/252 accessions demerged (2005 UniProt IDs resolved to current K-12 entries via REST API)",
    "HSP60: 325 raw -> 266 Tier-1 after excluding baits (2), co-chaperones (4), contaminants (4), SILAC filtering (median > 5)",
    "NDIC ('Not Detected In Control') = strongest enrichment evidence, imputed as 2x 95th percentile",
    "MitoCarta v2.0 -> v3.0: 70 reclassifications (52 respiratory chain subunits: matrix -> inner membrane)",
    "Compartment-matched controls: HSP60 vs matrix (525), not vs full human proteome (20,416)",
]
add_bullet_slide(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(2.2), items, font_size=13, spacing=Pt(4))

# ======================================================================
# SLIDE 7: CRITICAL METHODOLOGICAL DECISIONS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Nine Critical Methodological Decisions", "Each choice was made with explicit scientific rationale")
add_footer(slide, sn)

dec_data = [
    ["#", "Decision", "Rationale"],
    ["1", "CATH/Chainsaw for domain boundaries\n(NOT InterPro)", "InterPro = sequence domains (HMM). We need structural boundaries\nfor independently folding units. CATH = structural classification."],
    ["2", "Contact Order + FoldX for stability\n(NOT pLDDT alone)", "pLDDT = AlphaFold confidence, NOT thermodynamic stability.\nCO = folding kinetics. FoldX DeltaG = thermodynamics."],
    ["3", "OrthoFinder on full proteomes\n(RBH as supplementary)", "RBH misses many-to-many relationships. OrthoFinder captures\nparalogs. Increased pairs: 40 (RBH) -> 69 (merged)."],
    ["4", "SILAC-based HSP60 filtering", "Co-IP captures interactions, not function. SILAC ratio > 5\n+ MitoCarta confirmation = high-confidence substrates."],
    ["5", "Compartment + size-matched controls", "Controls for mitochondrial protein properties AND\nprotein size effects on domain distributions."],
    ["6", "Three-region N-terminal decomposition", "Pre-domain tail != transit peptide != first domain.\nConflating these produces meaningless results."],
    ["7", "MitoCarta 3.0 as ground truth", "Experimental evidence > predictors. 70 reclassifications\nbetween v2.0 and v3.0 significantly affect results."],
    ["8", "Hierarchical BH correction", "56 tests need correction. Hierarchical BH: correct within\nfamilies first, then across families. Controls FDR at 0.05."],
    ["9", "heiniglab/STRIDE (NOT bioconda)", "bioconda 'stride' = genomic variant caller.\nProtein STRIDE must be compiled from source."],
]
add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.8), dec_data,
          col_widths=[Inches(0.4), Inches(3.5), Inches(8.8)])

# ======================================================================
# SLIDE 8: OVERALL WORKFLOW
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Computational Workflow Overview", "Phase 1 (pilot, local Mac) + Phase 2 (full-scale, HPC)")
add_footer(slide, sn)

# Phase 1 box
add_shape(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.8), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.4),
            "Phase 1: Pilot (1,390 proteins, local Mac M1)", font_size=17, bold=True, color=MED_BLUE)

modules_p1 = [
    "Module A: Data Acquisition & Cleaning",
    "    - GroEL standardization (149 demerged accessions)",
    "    - HSP60 SILAC filtering (325 -> 266 Tier-1)",
    "    - MitoCarta 3.0 parsing",
    "Module B: Dataset Construction (FASTA extraction)",
    "Module C: Orthology & Homology",
    "    - RBH (40 pairs) + OrthoFinder (62 pairs) -> 69 merged",
    "Module D: Structure Acquisition & QC",
    "    - AlphaFold download (1,382 CIFs), DSSP, quality tiers",
    "Module E: Structural Domain Assignment",
    "    - CATH/Gene3D (1,151) + Chainsaw ML (236) = 99.8% coverage",
    "    - Foldseek clustering (1,155 clusters, 24 shared)",
    "Module F: N-vs-C Terminus Stability Analysis",
    "    - Three-region decomposition, contact order, pLDDT",
    "Module G: MTS Analysis (transit peptide-domain gap)",
    "Module H: Comparative Statistics (281 tests, 22 significant)",
    "Module I: Publication Figures (6 figures, PDF+PNG)",
]
add_bullet_slide(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.0), modules_p1, font_size=11, spacing=Pt(2))

# Phase 2 box
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.8), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
            "Phase 2: Full-Scale (25,007 proteins, HPC)", font_size=17, bold=True, color=ACCENT_ORANGE)

modules_p2 = [
    "Infrastructure: CSIR-IGIB HPC 'Tejas'",
    "    - SLURM scheduler, Lustre filesystem",
    "    - 40 CPUs, 380 GB RAM per node",
    "",
    "Step 1: AlphaFold bulk download (25,007 structures, ~22 GB)",
    "Steps 2-4: Foldseek pipeline",
    "    - createdb -> search (64 GB, 16 CPUs) -> cluster",
    "    - Result: 16,193 clusters, 27,063 proteins",
    "Step 5: Chainsaw full-scale (72h, 93.6% assigned)",
    "Step 6-7: FoldX thermodynamic stability",
    "    - 501 array tasks, ~40 sec/protein",
    "    - DeltaG = empirical free energy of folding",
    "    - IN PROGRESS (~42% complete)",
    "Step 8: Unified domain assignments (25,258 proteins)",
    "Step 9: Module F full-scale (2,648 paired comparisons)",
    "Step 10: Module H (56 tests, 25 significant)",
    "Step 11: Module I (6 polished publication figures)",
    "",
    "19 SLURM scripts with dependency chains",
    "All results transferred to local Mac",
]
add_bullet_slide(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(5.0), modules_p2, font_size=11, spacing=Pt(2))

# ======================================================================
# SLIDE 9: PHASE 2 HPC PIPELINE DIAGRAM
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Phase 2: HPC Pipeline Dependency Graph", "19 SLURM jobs with explicit dependency chains")
add_footer(slide, sn)

# Draw pipeline boxes
def pipeline_box(slide, left, top, width, height, text, color, font_size=10):
    shape = add_shape(slide, left, top, width, height, color)
    txBox = add_textbox(slide, left, top, width, height, text, font_size=font_size, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Vertically center
    return shape

# Row 1: Setup
pipeline_box(slide, Inches(5.5), Inches(1.3), Inches(2.5), Inches(0.5), "00: Environment Setup", MED_BLUE)

# Row 2: Download
pipeline_box(slide, Inches(5.5), Inches(2.1), Inches(2.5), Inches(0.5), "01: AlphaFold Download\n25,007 structures", MED_BLUE, 9)

# Row 3: Three branches
pipeline_box(slide, Inches(0.5), Inches(3.1), Inches(3.0), Inches(0.5), "02-04: Foldseek Pipeline\ncreatedb -> search -> cluster", ACCENT_GREEN, 9)
pipeline_box(slide, Inches(4.2), Inches(3.1), Inches(2.5), Inches(0.5), "05: Chainsaw\nML domains (72h)", LIGHT_BLUE, 9)
pipeline_box(slide, Inches(7.5), Inches(3.1), Inches(3.0), Inches(0.5), "06: FoldX Generate\n501 array tasks", ACCENT_ORANGE, 9)

# Row 4
pipeline_box(slide, Inches(0.5), Inches(4.1), Inches(3.0), Inches(0.5), "16,193 clusters\n27,063 proteins", ACCENT_GREEN, 9)
pipeline_box(slide, Inches(4.2), Inches(4.1), Inches(2.5), Inches(0.5), "08: Module E\nUnified domains", LIGHT_BLUE, 9)
pipeline_box(slide, Inches(7.5), Inches(4.1), Inches(3.0), Inches(0.5), "FoldX Array Job\n501 tasks (IN PROGRESS)", ACCENT_RED, 9)

# Row 5
pipeline_box(slide, Inches(7.5), Inches(5.0), Inches(3.0), Inches(0.5), "07: FoldX Collect\nMerge per-protein JSONs", ACCENT_ORANGE, 9)

# Row 6: Analysis chain
pipeline_box(slide, Inches(2.0), Inches(5.7), Inches(2.5), Inches(0.5), "09: Module F\nN-vs-C Stability", MED_BLUE, 9)
pipeline_box(slide, Inches(5.2), Inches(5.7), Inches(2.5), Inches(0.5), "10: Module H\nStatistics (56 tests)", MED_BLUE, 9)
pipeline_box(slide, Inches(8.4), Inches(5.7), Inches(2.5), Inches(0.5), "11: Module I\n6 Publication Figures", MED_BLUE, 9)

# Resource annotation
add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
            "Resources: Foldseek search = 64 GB RAM, 16 CPUs | Chainsaw = 72h wall time | FoldX = 501 array tasks x ~40s/protein | Total data: ~22 GB structures",
            font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 10: STRUCTURAL DOMAIN ASSIGNMENT METHOD
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Method: Structural Domain Assignment", "Module E: CATH + Chainsaw ML = 99.8% coverage")
add_footer(slide, sn)

# CATH box
add_shape(slide, Inches(0.4), Inches(1.4), Inches(4.0), Inches(3.0), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(3.6), Inches(0.4),
            "CATH / Gene3D (Primary)", font_size=16, bold=True, color=MED_BLUE)
items = [
    "Curated structural domain database",
    "4-level hierarchy: Class.Architecture.Topology.Homology",
    "Queried via InterPro REST API",
    "1,151/1,390 pilot proteins (82.8%)",
    "Phase 2: 1,390 proteins with CATH annotations",
    "Mean domain pLDDT: 92.1",
]
add_bullet_slide(slide, Inches(0.6), Inches(2.0), Inches(3.6), Inches(2.2), items, font_size=12, spacing=Pt(3))

# Chainsaw box
add_shape(slide, Inches(4.7), Inches(1.4), Inches(4.0), Inches(3.0), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(4.9), Inches(1.5), Inches(3.6), Inches(0.4),
            "Chainsaw ML (Secondary)", font_size=16, bold=True, color=ACCENT_ORANGE)
items = [
    "Deep learning model trained on CATH domains",
    "Predicts structural boundaries from 3D coords",
    "Uses STRIDE secondary structure as input",
    "236 additional pilot proteins assigned",
    "Phase 2: 23,868 proteins (full proteomes)",
    "Unified: 25,258 total domain records",
]
add_bullet_slide(slide, Inches(4.9), Inches(2.0), Inches(3.6), Inches(2.2), items, font_size=12, spacing=Pt(3))

# Foldseek box
add_shape(slide, Inches(9.0), Inches(1.4), Inches(4.0), Inches(3.0), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(9.2), Inches(1.5), Inches(3.6), Inches(0.4),
            "Foldseek Clustering", font_size=16, bold=True, color=ACCENT_GREEN)
items = [
    "3Di + AA structural alphabet",
    "All-vs-all structural comparison",
    "BLAST-like speed for 3D search",
    "Pilot: 1,155 clusters, 24 shared",
    "Phase 2: 16,193 clusters, 27,063 proteins",
    "75.6% singletons (unique folds)",
]
add_bullet_slide(slide, Inches(9.2), Inches(2.0), Inches(3.6), Inches(2.2), items, font_size=12, spacing=Pt(3))

# Coverage summary
cov_data = [
    ["Scale", "CATH", "Chainsaw", "Total Coverage"],
    ["Phase 1 (pilot)", "1,151 (82.8%)", "236 (17.0%)", "1,387/1,390 (99.8%)"],
    ["Phase 2 (full)", "1,390", "23,868", "25,258 proteins"],
]
add_table(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(0.9), cov_data)

# CATH hierarchy
add_textbox(slide, Inches(0.6), Inches(5.9), Inches(12), Inches(0.3),
            "CATH Hierarchy Example: 3.20.20.70 = Alpha-Beta (3) | Barrel (20) | TIM Barrel (20) | Aldolase Class I (70)",
            font_size=13, bold=True, color=DARK_BLUE)

# ======================================================================
# SLIDE 11: CONTACT ORDER & STABILITY METHOD
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Method: Stability Metrics", "Contact Order (kinetics) + FoldX DeltaG (thermodynamics)")
add_footer(slide, sn)

# Contact order
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.8), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(0.4),
            "Relative Contact Order (RCO)", font_size=17, bold=True, color=MED_BLUE)
items = [
    "Plaxco, Simons & Baker (1998) definition",
    "RCO = (1/N_contacts x L) x SUM |i - j|",
    "Where: i,j = CA atoms within 8 A, |i-j| >= 6 residues",
    "High RCO = topologically complex, folds SLOWLY",
    "Correlates with experimental folding rate (r ~ -0.75)",
    "Computed per region: N-domain vs C-region",
]
add_bullet_slide(slide, Inches(0.6), Inches(2.0), Inches(5.6), Inches(2.0), items, font_size=13, spacing=Pt(3))

# FoldX
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.8), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4),
            "FoldX Thermodynamic Stability (DeltaG)", font_size=17, bold=True, color=ACCENT_ORANGE)
items = [
    "Semi-empirical force field (FoldX 5.1)",
    "Accounts for: van der Waals, solvation, H-bonds,",
    "    electrostatics, entropy",
    "More negative DeltaG = more stable fold",
    "~40 seconds per protein",
    "STATUS: 42% complete on HPC (501 array tasks)",
]
add_bullet_slide(slide, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0), items, font_size=13, spacing=Pt(3))

# Three region decomposition
add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
            "Three-Region Decomposition (per protein)", font_size=17, bold=True, color=DARK_BLUE)

# Visual diagram of three regions
add_shape(slide, Inches(1.0), Inches(5.2), Inches(2.0), Inches(0.7), RGBColor(0xDD, 0xDD, 0xDD))
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(2.0), Inches(0.5),
            "Pre-domain Tail\n(may include MTS)", font_size=10, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(3.2), Inches(5.2), Inches(3.5), Inches(0.7), MED_BLUE)
add_textbox(slide, Inches(3.2), Inches(5.3), Inches(3.5), Inches(0.5),
            "N-Domain\n(First structural domain)", font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(6.9), Inches(5.2), Inches(5.0), Inches(0.7), ACCENT_ORANGE)
add_textbox(slide, Inches(6.9), Inches(5.3), Inches(5.0), Inches(0.5),
            "C-Region\n(All subsequent domains + linkers)", font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.0), Inches(6.1), Inches(1.5), Inches(0.3), "Residue 1", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.2), Inches(6.1), Inches(3.5), Inches(0.3), "Domain start --- Domain end", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.5), Inches(6.1), Inches(2.5), Inches(0.3), "Protein end", font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

# pLDDT note
add_textbox(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.3),
            "Important: pLDDT is AlphaFold's confidence metric, NOT stability. We report it alongside CO and DeltaG but do NOT equate it with thermodynamic stability.",
            font_size=11, bold=False, color=ACCENT_RED)

# ======================================================================
# SLIDE 12: RESULTS — DOMAIN ARCHITECTURE (with figure)
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results: Goal 1 — Domain Architecture", "Chaperonin substrates enriched for specific complex folds")
add_footer(slide, sn)

# Add figure
fig1_path = os.path.join(FIG_DIR, "fig1_domain_architecture.png")
if os.path.exists(fig1_path):
    slide.shapes.add_picture(fig1_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.5))

# Key findings below figure
add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.3),
            "Key Findings (Phase 2, 25,007 proteins)", font_size=16, bold=True, color=DARK_BLUE)

items = [
    "GroEL enriched: TIM barrel (3.20.20.70) OR=8.4, p=1.9x10^-8  |  Winged helix (1.10.10.10) OR=50.9, p=3.5x10^-9",
    "HSP60 enriched: 3.30.830.10 OR=5.4, p=2.5x10^-4  |  3.90.226.10 OR=4.8, p=4.6x10^-4",
    "Alpha-beta class dominates all groups (60-71%). NOT enriched for multi-domain proteins (OR=1.13, NS)",
    "Interpretation: Chaperonin substrate identity = specific fold topology (TIM barrels), NOT general complexity",
]
add_bullet_slide(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8), items, font_size=13, spacing=Pt(4))

# ======================================================================
# SLIDE 13: DOMAIN ENRICHMENT TABLE
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Domain Architecture: Enrichment Details", "Fisher's exact test with size-matched controls")
add_footer(slide, sn)

# GroEL enrichments
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
            "GroEL Substrate Enrichments", font_size=16, bold=True, color=MED_BLUE)

groel_data = [
    ["CATH Code", "Name", "OR", "p (BH)", "Interp."],
    ["3.20.20.70", "TIM barrel / Aldolase I", "8.4", "1.9x10^-8", "Complex 8-fold barrel topology"],
    ["1.10.10.10", "Winged helix-like", "50.9", "3.5x10^-9", "DNA-binding, complex HTH"],
    ["3.30.420.40", "Nucleotidyltransferase", "6.0", "5.8x10^-3", "Enzyme fold"],
    ["3.40.640.10", "Muconolactone isomerase", "2.6", "1.5x10^-2", "Metabolic enzyme"],
]
add_table(slide, Inches(0.3), Inches(1.7), Inches(6.2), Inches(1.8), groel_data)

# HSP60 enrichments
add_textbox(slide, Inches(0.5), Inches(3.8), Inches(5), Inches(0.3),
            "HSP60 Substrate Enrichments", font_size=16, bold=True, color=ACCENT_ORANGE)

hsp60_data = [
    ["CATH Code", "Name", "OR", "p (BH)", "Interp."],
    ["3.30.830.10", "Rossmann-like", "5.4", "2.5x10^-4", "Matrix enzyme fold"],
    ["3.90.226.10", "-", "4.8", "4.6x10^-4", "Mitochondrial fold"],
    ["3.40.50.620", "-", "3.3", "1.3x10^-2", "Alpha-beta fold"],
    ["2.40.30.10", "-", "3.6", "1.2x10^-2", "Beta sandwich"],
]
add_table(slide, Inches(0.3), Inches(4.2), Inches(6.2), Inches(1.8), hsp60_data)

# Right side: Why TIM barrels
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.0), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
            "Why TIM Barrels Need Chaperonins", font_size=16, bold=True, color=ACCENT_ORANGE)
items = [
    "TIM barrel = 8-stranded beta/alpha barrel",
    "Complex topology requiring 8 beta-strands and",
    "    8 alpha-helices in precise alternating arrangement",
    "",
    "The barrel CANNOT form until all 8 strands are present",
    "Co-translational folding is IMPOSSIBLE",
    "",
    "Protein must be fully synthesized before folding ->",
    "    Creates a window of vulnerability",
    "    Unfolded protein aggregates without protection",
    "",
    "GroEL barrel provides protected folding environment",
    "OR = 8.4 means TIM barrel substrates are 8.4x",
    "    more common in GroEL than expected by chance",
]
add_bullet_slide(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(4.2), items, font_size=12, spacing=Pt(3))

# ======================================================================
# SLIDE 14: RESULTS — N-vs-C STABILITY (with figure)
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results: Goal 2 — N-vs-C Contact Order Asymmetry", "N-domains have HIGHER contact order than C-regions — UNIVERSALLY")
add_footer(slide, sn)

fig2_path = os.path.join(FIG_DIR, "fig2_n_vs_c_stability.png")
if os.path.exists(fig2_path):
    slide.shapes.add_picture(fig2_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.2))

# Results table
rco_data = [
    ["Dataset", "n pairs", "Median N-C diff", "Effect size r", "p-value", "Significant?"],
    ["GroEL substrates", "124", "0.043", "0.41", "8.9x10^-5", "YES"],
    ["HSP60 substrates", "131", "0.059", "0.46", "5.3x10^-6", "YES"],
    ["Matrix background", "251", "0.069", "0.43", "2.4x10^-9", "YES"],
    ["Mito background", "425", "0.064", "0.48", "7.1x10^-18", "YES (strongest)"],
]
add_table(slide, Inches(0.3), Inches(4.7), Inches(7.5), Inches(1.8), rco_data)

# Interpretation box
add_shape(slide, Inches(8.2), Inches(4.7), Inches(4.8), Inches(2.3), RGBColor(0xFC, 0xE4, 0xE4))
add_textbox(slide, Inches(8.4), Inches(4.8), Inches(4.4), Inches(0.3),
            "THE SURPRISE", font_size=16, bold=True, color=ACCENT_RED)
items = [
    "N > C contact order is UNIVERSAL",
    "Present in ALL datasets, not just substrates",
    "Effect is medium-to-large (r = 0.41-0.48)",
    "Most significant in general mito background!",
    "",
    "This is NOT a chaperonin-specific feature",
    "It reflects physics of vectorial translation",
]
add_bullet_slide(slide, Inches(8.4), Inches(5.2), Inches(4.4), Inches(1.6), items, font_size=12, spacing=Pt(2))

# ======================================================================
# SLIDE 15: THE CRUCIAL NEGATIVE RESULTS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "The Crucial Negative Results", "What we did NOT find is as important as what we found")
add_footer(slide, sn)

# H2.2 rejection
add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.5), RGBColor(0xFC, 0xE4, 0xE4))
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.6), Inches(0.4),
            "H2.2 REJECTED: Asymmetry NOT Substrate-Specific", font_size=15, bold=True, color=ACCENT_RED)
items = [
    "Mann-Whitney U: substrates vs background",
    "GroEL vs E. coli bg: p = 0.058 (NOT significant)",
    "HSP60 vs mito bg: p = 0.536 (NOT significant)",
    "",
    "Background proteins show the SAME asymmetry",
    "The effect is universal to all multi-domain proteins",
]
add_bullet_slide(slide, Inches(0.6), Inches(2.0), Inches(5.6), Inches(1.8), items, font_size=13, spacing=Pt(3))

# H2.3 rejection
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.5), RGBColor(0xFC, 0xE4, 0xE4))
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4),
            "H2.3 REJECTED: No GroEL Class Gradient", font_size=15, bold=True, color=ACCENT_RED)
items = [
    "Kruskal-Wallis H-test across Class I/II/III:",
    "Contact order: p = 0.77",
    "pLDDT: p = 0.92",
    "",
    "If asymmetry drove dependence, Class III should",
    "show greatest effect. They don't.",
]
add_bullet_slide(slide, Inches(7.0), Inches(2.0), Inches(5.6), Inches(1.8), items, font_size=13, spacing=Pt(3))

# GroEL class figure
fig3_path = os.path.join(FIG_DIR, "fig3_groel_class_comparison.png")
if os.path.exists(fig3_path):
    slide.shapes.add_picture(fig3_path, Inches(0.3), Inches(4.2), Inches(12.7), Inches(2.8))

# ======================================================================
# SLIDE 16: BIOLOGICAL INTERPRETATION OF N-vs-C
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Biological Interpretation: Why N > C Contact Order?", "A 'gravitational constant' of protein architecture")
add_footer(slide, sn)

items = [
    "1. Proteins synthesized N-to-C on the ribosome",
    "",
    "2. N-terminal region emerges FIRST, begins to fold while rest is still being synthesized",
    "",
    "3. Evolutionary pressure places topologically complex folds (high CO) at N-terminus,",
    "     where they have the MOST TIME to fold during translation",
    "",
    "4. C-terminal regions, emerging LAST, tend to adopt simpler folds or extend existing domains",
    "",
    "5. This is a FUNDAMENTAL property of protein organization:",
    "     - Conserved across 2 billion years of evolution",
    "     - Present in ALL multi-domain proteins",
    "     - NOT specific to chaperonin substrates",
    "     - A 'gravitational constant' — always present, not caused by chaperonins",
    "",
    "Key insight: Chaperonins don't preferentially assist C-terminal folding.",
    "They capture the ENTIRE protein and help the globally complex topology snap into place.",
]
add_bullet_slide(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5), items, font_size=15, spacing=Pt(4))

# ======================================================================
# SLIDE 17: RESULTS — MTS TARGETING (with figure)
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results: Goal 3 — Mitochondrial Targeting Signal Architecture", "HSP60 substrates are matrix-enriched; MTS = pre-domain extension")
add_footer(slide, sn)

fig4_path = os.path.join(FIG_DIR, "fig4_mts_targeting.png")
if os.path.exists(fig4_path):
    slide.shapes.add_picture(fig4_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.2))

# Results
add_textbox(slide, Inches(0.5), Inches(4.7), Inches(6), Inches(0.3),
            "H3.1: Matrix Enrichment", font_size=16, bold=True, color=ACCENT_GREEN)
items = [
    "HSP60 substrates: 46.6% matrix localization",
    "Background mitochondrial: 28.5% matrix",
    "OR = 3.29, 95% CI [2.46, 4.40], p = 1.6x10^-16",
    "21.1% non-canonical matrix import (no MTS detected)",
]
add_bullet_slide(slide, Inches(0.5), Inches(5.1), Inches(6), Inches(1.3), items, font_size=13, spacing=Pt(3))

add_textbox(slide, Inches(6.8), Inches(4.7), Inches(6), Inches(0.3),
            "H3.3: MTS = Pre-Domain Extension", font_size=16, bold=True, color=ACCENT_GREEN)
items = [
    "84.4% of transit peptides end BEFORE first domain",
    "Only 15.6% overlap with structural domain",
    "Binomial test: p = 3.4x10^-51",
    "Median gap: 18 residues (MTS cleavage to domain start)",
]
add_bullet_slide(slide, Inches(6.8), Inches(5.1), Inches(6), Inches(1.3), items, font_size=13, spacing=Pt(3))

# Landing pad model
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5),
            '"Landing Pad" Model: After import & MTS cleavage, protein emerges with short unstructured linker (~18 aa) followed by unfolded first domain — optimally positioned for HSP60 capture and assisted folding.',
            font_size=13, bold=True, color=ACCENT_GREEN)

# ======================================================================
# SLIDE 18: CROSS-SPECIES CONSERVATION (with figure)
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results: Cross-Species Conservation", "N-domain complexity conserved across 2 billion years (r = 0.84)")
add_footer(slide, sn)

fig5_path = os.path.join(FIG_DIR, "fig5_orthology.png")
if os.path.exists(fig5_path):
    slide.shapes.add_picture(fig5_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.3))

items = [
    "69 cross-species homolog pairs (33 both methods, 7 RBH-only, 29 orthogroup-only)",
    "",
    "N-domain contact order: Spearman r = 0.84, p = 5.3x10^-13 (n = 45 paired values)",
    "    -> Topological complexity of N-terminal domain is STRONGLY conserved across species",
    "",
    "Fold type conservation: 55/69 pairs (79.7%) share same top CATH superfamily",
    "    -> Chaperonin substrates constrained to same structural fold across species",
    "",
    "RBH observation: Only 8/84 (9.5%) Class III obligate substrates have HSP60 orthologs",
    "    vs 18.4% of Class I, 19.0% of Class II",
    "    -> Most GroEL-dependent proteins are the most divergent from human counterparts",
]
add_bullet_slide(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(2.5), items, font_size=14, spacing=Pt(4))

# ======================================================================
# SLIDE 19: SUMMARY FIGURE
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results Summary", "25 of 56 tests significant after hierarchical BH correction")
add_footer(slide, sn)

fig6_path = os.path.join(FIG_DIR, "fig6_summary.png")
if os.path.exists(fig6_path):
    slide.shapes.add_picture(fig6_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.5))

# Statistical summary table
stat_data = [
    ["Family", "Tests", "Significant", "Key Result"],
    ["Domain Architecture", "24", "9 (37.5%)", "TIM barrel OR=8.4; 1.10.10.10 OR=50.9"],
    ["N-vs-C Asymmetry", "30", "14 (46.7%)", "N>C universal (r=0.41-0.48); substrate NS"],
    ["MTS Targeting", "2", "2 (100%)", "Matrix OR=3.29; pre-domain 84.4%"],
    ["TOTAL", "56", "25 (44.6%)", "Hierarchical BH, FDR < 0.05"],
]
add_table(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), stat_data)

# ======================================================================
# SLIDE 20: BIOLOGICAL SYNTHESIS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Biological Synthesis", "The three goals converge into a coherent narrative")
add_footer(slide, sn)

# Three boxes
add_shape(slide, Inches(0.3), Inches(1.4), Inches(4.1), Inches(2.8), VERY_LIGHT_BLUE)
add_textbox(slide, Inches(0.5), Inches(1.5), Inches(3.7), Inches(0.4),
            "What Makes a Chaperonin\nSubstrate?", font_size=14, bold=True, color=MED_BLUE)
items = [
    "It's the FOLD, not general difficulty",
    "TIM barrels: can't fold co-translationally",
    "8-fold topology needs all elements present",
    "NOT about domain count or size",
    "NOT about N-vs-C asymmetry",
    "NOT about obligate dependence (class effect)",
]
add_bullet_slide(slide, Inches(0.5), Inches(2.1), Inches(3.7), Inches(1.9), items, font_size=11, spacing=Pt(2))

add_shape(slide, Inches(4.6), Inches(1.4), Inches(4.1), Inches(2.8), RGBColor(0xFD, 0xF2, 0xE0))
add_textbox(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.4),
            "How Does the Mito\nSystem Work?", font_size=14, bold=True, color=ACCENT_ORANGE)
items = [
    "HSP60 substrates = matrix residents (3.3x)",
    "MTS creates a 'landing pad'",
    "84.4% pre-domain extension",
    "~18 residue gap after cleavage",
    "First domain emerges unfolded",
    "Optimal for HSP60 capture",
]
add_bullet_slide(slide, Inches(4.8), Inches(2.1), Inches(3.7), Inches(1.9), items, font_size=11, spacing=Pt(2))

add_shape(slide, Inches(8.9), Inches(1.4), Inches(4.1), Inches(2.8), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(9.1), Inches(1.5), Inches(3.7), Inches(0.4),
            "What is Conserved\nAcross 2 Billion Years?", font_size=14, bold=True, color=ACCENT_GREEN)
items = [
    "The fold itself (79.7% same CATH SF)",
    "N-domain complexity (r = 0.84)",
    "TIM barrels need GroEL in E. coli,",
    "    HSP60 in humans",
    "N>C asymmetry = universal property",
    "Translational history conserved",
]
add_bullet_slide(slide, Inches(9.1), Inches(2.1), Inches(3.7), Inches(1.9), items, font_size=11, spacing=Pt(2))

# The title explained
add_shape(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.6), DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
            '"The End is the Beginning" — Explained', font_size=18, bold=True, color=ACCENT_ORANGE)
items = [
    "The C-terminus ('end' of translation) is the 'beginning' of chaperonin engagement",
    "BUT — not because C-terminal regions are more complex",
    "The entire protein, synthesized N-to-C, has its most complex domain at the N-terminus",
    "The chaperonin captures the entire protein and helps the globally complex topology form",
    "The N>C asymmetry is a footprint of translational history, conserved across all of life",
    "Chaperonin substrate identity is determined by FOLD TOPOLOGY, not positional complexity",
]
add_bullet_slide(slide, Inches(0.6), Inches(5.1), Inches(12), Inches(1.8), items, font_size=13, color=WHITE, spacing=Pt(3))

# ======================================================================
# SLIDE 21: LIMITATIONS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Limitations and Caveats")
add_footer(slide, sn)

# Methodological
add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.3),
            "Methodological", font_size=16, bold=True, color=ACCENT_RED)
items = [
    "pLDDT = confidence, not stability (mitigated by CO + FoldX)",
    "Co-IP captures interactions, not function (mitigated by SILAC)",
    "AlphaFold structures are predictions (mean pLDDT 85-93)",
    "Contact order imperfect for multi-state folders",
    "MitoCarta annotations may evolve in future versions",
]
add_bullet_slide(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(2.2), items, font_size=13, spacing=Pt(3))

# Statistical
add_textbox(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.3),
            "Statistical", font_size=16, bold=True, color=ACCENT_RED)
items = [
    "56 tests: risk of false positives despite BH correction",
    "Only 69 homolog pairs — limited power for cross-species",
    "10 kDa size-matching bins are approximate",
    "No Class III equivalents for HSP60 substrates",
]
add_bullet_slide(slide, Inches(6.8), Inches(1.7), Inches(6), Inches(2.0), items, font_size=13, spacing=Pt(3))

# Biological
add_textbox(slide, Inches(0.5), Inches(4.0), Inches(6), Inches(0.3),
            "Biological", font_size=16, bold=True, color=ACCENT_RED)
items = [
    "~30-40% of matrix proteins lack detectable MTS",
    "No TargetP/SignalP6 predictions (DTU license required)",
    "HSP60 lacks dependence class annotations",
    "Alternative import pathways not characterized",
]
add_bullet_slide(slide, Inches(0.5), Inches(4.4), Inches(6), Inches(1.8), items, font_size=13, spacing=Pt(3))

# Pending
add_textbox(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(0.3),
            "Pending / In Progress", font_size=16, bold=True, color=ACCENT_ORANGE)
items = [
    "FoldX DeltaG: ~42% complete on HPC",
    "    -> Will add thermodynamic stability dimension",
    "    -> Re-run Modules F, H, I after completion",
    "    -> Estimated completion: April 1-2, 2026",
    "Manuscript preparation after FoldX integration",
]
add_bullet_slide(slide, Inches(6.8), Inches(4.4), Inches(6), Inches(1.8), items, font_size=13, spacing=Pt(3))

# ======================================================================
# SLIDE 22: FUTURE WORK
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Future Directions", "After FoldX completion and manuscript preparation")
add_footer(slide, sn)

items = [
    "Immediate (Session 7, ~April 2026):",
    "    1. Verify FoldX completion (25,007 proteins, check success rate)",
    "    2. Re-run Module F with FoldX DeltaG integration (N-vs-C thermodynamic stability)",
    "    3. Re-run Module H with DeltaG-enriched statistical tests",
    "    4. Re-run Module I for updated publication figures with DeltaG violins",
    "    5. Transfer final results to local Mac",
    "    6. Manuscript preparation",
    "",
    "Longer-term:",
    "    7. TargetP 2.0 predictions for MTS gap characterization (requires DTU license)",
    "    8. IUPred2A for intrinsic disorder prediction in pre-domain tails",
    "    9. Molecular dynamics validation of FoldX stability predictions",
    "    10. Extend to Group II chaperonins (TRiC/CCT in eukaryotic cytoplasm)",
    "    11. Deep learning models for chaperonin substrate prediction (features: fold type, CO, DeltaG)",
    "    12. Experimental validation of 'landing pad' model via truncation experiments",
]
add_bullet_slide(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.5), items, font_size=15, spacing=Pt(5), bold_prefix=True)

# ======================================================================
# SLIDE 23: SOFTWARE & REPRODUCIBILITY
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Software, Data, and Reproducibility")
add_footer(slide, sn)

sw_data = [
    ["Tool", "Version", "Purpose"],
    ["Python", "3.9.16 (Anaconda)", "Core analysis language"],
    ["pandas", "2.2.2", "Data manipulation"],
    ["scipy", "1.9.2", "Statistical tests"],
    ["BioPython", "1.78", "Sequence/structure parsing"],
    ["MMseqs2", "18.8cc5c", "Sequence search, RBH, orthology"],
    ["Foldseek", "10.941cd33", "Structural similarity search"],
    ["Chainsaw", "2024 (latest)", "ML domain boundary prediction"],
    ["STRIDE", "heiniglab/stride", "Secondary structure assignment"],
    ["FoldX", "5.1 (2027 build)", "Thermodynamic stability (DeltaG)"],
    ["mkDSSP", "2.2.1", "DSSP secondary structure"],
    ["Snakemake", "7.32.4", "Workflow management"],
    ["AlphaFold DB", "v4/v6", "Predicted 3D structures (CIF)"],
]
add_table(slide, Inches(0.3), Inches(1.3), Inches(6.3), Inches(4.5), sw_data)

# Reproducibility
add_textbox(slide, Inches(7.0), Inches(1.3), Inches(6), Inches(0.3),
            "Reproducibility", font_size=16, bold=True, color=DARK_BLUE)
items = [
    "All code: workflow/scripts/ (Phase 1) + workflow/phase2/ (Phase 2)",
    "All raw data downloadable from public databases",
    "Phase 2 deployment guide: docs/HPC_DEPLOYMENT_GUIDE.md",
    "Central config: workflow/phase2/config.yaml",
    "19 SLURM scripts with resource specifications",
    "Snakemake workflow for automated orchestration",
    "",
    "Data scale:",
    "    Phase 1: 1,390 proteins, ~466 MB structures",
    "    Phase 2: 25,007 proteins, ~22 GB structures",
    "    Total results: ~186 MB",
    "    6 publication figures (PDF + PNG, 300 DPI)",
]
add_bullet_slide(slide, Inches(7.0), Inches(1.7), Inches(6), Inches(4.5), items, font_size=13, spacing=Pt(4))

# ======================================================================
# SLIDE 24: CONCLUSIONS
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
            "Conclusions", font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Key conclusions
conclusions = [
    "1. Chaperonin substrates are enriched for SPECIFIC complex fold topologies",
    "      GroEL: TIM barrels (OR=8.4), Winged helix (OR=50.9)",
    "      HSP60: Rossmann-like (OR=5.4), other matrix enzyme folds",
    "",
    "2. N-terminal domains have higher contact order than C-terminal regions",
    "      Effect is UNIVERSAL — present in ALL multi-domain proteins",
    "      NOT substrate-specific (Mann-Whitney substrate vs bg: NS)",
    "      NO GroEL class gradient (Kruskal-Wallis p=0.77)",
    "",
    "3. Mitochondrial targeting signals are predominantly pre-domain extensions",
    "      84.4% pre-domain (p = 3.4x10^-51), median gap = 18 residues",
    "      HSP60 substrates 3.3x enriched for matrix (p = 1.6x10^-16)",
    "      The 'landing pad' model for post-import chaperonin engagement",
    "",
    "4. Chaperonin substrate properties are CONSERVED across 2 billion years",
    "      N-domain contact order: r = 0.84 (p = 5.3x10^-13)",
    "      79.7% of homolog pairs share the same CATH superfamily",
    "",
    "5. The N>C asymmetry is a 'gravitational constant' of protein architecture",
    "      Reflects co-translational folding physics, not chaperonin biology",
]
add_bullet_slide(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(5.8), conclusions, font_size=15, color=WHITE, spacing=Pt(3))

# ======================================================================
# SLIDE 25: THANK YOU
# ======================================================================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

add_shape(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT_ORANGE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(0.7),
            "Thank You", font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(0.5),
            "Antah Asti Prarambh — The End is the Beginning", font_size=20, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11.333), Inches(0.4),
            "Vishal Bharti", font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.0), Inches(11.333), Inches(0.4),
            "CSIR-Institute of Genomics and Integrative Biology (IGIB), New Delhi", font_size=15, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(1.5),
            "Pipeline: 16 Python modules + 19 SLURM scripts\n"
            "Data: 25,007 proteins, 7 datasets, 56 statistical tests\n"
            "Figures: 6 publication-quality figures (PDF + PNG)\n"
            "Status: Phase 2 complete; FoldX ~42% done on HPC",
            font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

# Save
prs.save(OUT_FILE)
print(f"Presentation saved to: {OUT_FILE}")
print(f"Total slides: {sn}")
