#!/usr/bin/env python3
"""
Comprehensive PowerPoint presentation for Antah Asti Prarambh project.
V2: Full technical details, exact parameters, tool specifications, rationale.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE, "results/phase2/figures")
OUT_FILE = os.path.join(BASE, "Antah_Asti_Prarambh_Presentation.pptx")

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x86, 0xC8)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_GREEN = RGBColor(0x2D, 0x8E, 0x4E)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
VL_BLUE = RGBColor(0xE8, 0xF0, 0xF8)
VL_ORANGE = RGBColor(0xFD, 0xF2, 0xE0)
VL_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
VL_RED = RGBColor(0xFC, 0xE4, 0xE4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL_SLIDES = 35

def add_bg(slide, color=WHITE):
    bg = slide.background; f = bg.fill; f.solid(); f.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, name="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = name; p.alignment = align
    return box

def bullets(slide, l, t, w, h, items, sz=14, color=BLACK, sp=Pt(4), bp=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = sp
        if bp and ": " in item and not item.startswith(" "):
            pre, rest = item.split(": ", 1)
            r1 = p.add_run(); r1.text = pre + ": "; r1.font.size = Pt(sz); r1.font.bold = True; r1.font.color.rgb = color; r1.font.name = "Calibri"
            r2 = p.add_run(); r2.text = rest; r2.font.size = Pt(sz); r2.font.bold = False; r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = item; r.font.size = Pt(sz); r.font.color.rgb = color; r.font.name = "Calibri"
    return box

def header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BLUE)
    tb(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.55), title, sz=28, bold=True, color=WHITE)
    if subtitle:
        tb(slide, Inches(0.6), Inches(0.62), Inches(12), Inches(0.4), subtitle, sz=14, color=RGBColor(0xBB,0xCC,0xDD))

def footer(slide, n):
    tb(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3), f"{n}/{TOTAL_SLIDES}", sz=10, color=GRAY, align=PP_ALIGN.RIGHT)
    tb(slide, Inches(0.5), Inches(7.1), Inches(5), Inches(0.3), "Antah Asti Prarambh | Vishal Bharti, CSIR-IGIB", sz=10, color=GRAY)

def table(slide, l, t, w, h, data, cw=None, hc=DARK_BLUE):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, l, t, w, h); tbl = ts.table
    if cw:
        for i, ww in enumerate(cw): tbl.columns[i].width = ww
    for r, rd in enumerate(data):
        for c, ct in enumerate(rd):
            cell = tbl.cell(r, c); cell.text = str(ct)
            for pg in cell.text_frame.paragraphs:
                pg.font.size = Pt(10); pg.font.name = "Calibri"
                if r == 0: pg.font.bold = True; pg.font.color.rgb = WHITE; pg.alignment = PP_ALIGN.CENTER
                else: pg.font.color.rgb = BLACK
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = hc
            elif r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = VL_BLUE
    return ts

sn = 0

# ===================== SLIDE 1: TITLE =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
tb(slide, Inches(1), Inches(0.8), Inches(11.3), Inches(0.7), "ANTAH ASTI PRARAMBH", sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(1.7), Inches(11.3), Inches(0.5), '"The End is the Beginning"', sz=24, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(0.8), "A Comparative Structural Proteomics Study of\nGroup I Chaperonin Substrates", sz=26, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5), "GroEL/GroES (E. coli)  vs  HSP60/HSP10 (Human Mitochondria)", sz=18, color=RGBColor(0xAA,0xCC,0xEE), align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.4), "Vishal Bharti", sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4), "CSIR-Institute of Genomics and Integrative Biology (IGIB), New Delhi", sz=14, color=RGBColor(0xBB,0xCC,0xDD), align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.4), "March 2026", sz=14, color=RGBColor(0xBB,0xCC,0xDD), align=PP_ALIGN.CENTER)

# ===================== SLIDE 2: OUTLINE =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE); header(slide, "Presentation Outline"); footer(slide, sn)
left = ["1.  Background & Motivation", "2.  Three Scientific Goals & Pre-Registered Hypotheses", "3.  Seven Datasets & Data Cleaning",
        "4.  Critical Methodological Decisions (9 decisions)", "5.  Tools & Software Stack",
        "6.  Phase 1 Workflow: Modules A-I (Pilot, 1,390 proteins)",
        "7.  Module C: Orthology — MMseqs2 Parameters & OrthoFinder",
        "8.  Module D: AlphaFold & DSSP — Exact Parameters",
        "9.  Module E: CATH + Chainsaw + Foldseek — Parameters"]
right = ["10. Module F: Contact Order — Plaxco Definition & Parameters",
         "11. Module G: MTS Analysis — UniProt API & Gap Calculation",
         "12. Module H: Statistics — Hierarchical BH Framework",
         "13. Phase 2: HPC Pipeline — SLURM Resource Allocation",
         "14. FoldX Thermodynamic Stability — Parameters",
         "15. Results: Domain Architecture (Figures)",
         "16. Results: N-vs-C Asymmetry (Figures)",
         "17. Results: MTS Targeting & Cross-Species Conservation",
         "18. Biological Synthesis, Limitations, Conclusions"]
bullets(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(6.0), left, sz=14, sp=Pt(6))
bullets(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(6.0), right, sz=14, sp=Pt(6))

# ===================== SLIDE 3: BACKGROUND =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Background: Chaperonins in Protein Folding", "Barrel-shaped complexes that assist protein folding (~10-15% of proteome)"); footer(slide, sn)
add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.0), Inches(5.5), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(5.6), Inches(0.4), "GroEL/GroES (E. coli)", sz=20, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.9), Inches(5.6), Inches(4.8), [
    "~800 kDa double-ring complex (2x7 subunits)",
    "252 substrate proteins (Kerner et al., Cell 2005, Table S3)",
    "Three dependency classes:",
    "    Class I (38): Spontaneous folders, accelerated by GroEL",
    "    Class II (126): Partially dependent, fold slowly without GroEL",
    "    Class III (84): OBLIGATE — cannot fold, aggregate without GroEL",
    "Co-chaperonin: GroES (10 kDa lid, 7-mer)",
    "Encapsulates unfolded substrate in ~85 A cavity",
    "ATP-driven conformational cycle (~10 sec per round)",
], sz=12, sp=Pt(3))
add_rect(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5), VL_ORANGE)
tb(slide, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.4), "HSP60/HSP10 (Human Mitochondria)", sz=20, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(7.2), Inches(1.9), Inches(5.6), Inches(4.8), [
    "Mitochondrial matrix chaperonin (HSPD1 gene)",
    "266 Tier-1 substrates (Morten et al., Mol Cell 2020)",
    "    Identified by SILAC-quantified co-immunoprecipitation",
    "Nuclear-encoded, imported through TOM/TIM complexes",
    "Must UNFOLD for import, then RE-FOLD in matrix",
    "Co-chaperonin: HSP10 (HSPE1, 10 kDa, 7-mer lid)",
    "Evolutionary homolog of GroEL",
    "    Diverged ~2 billion years ago (endosymbiosis)",
    "    Alpha-proteobacterial ancestor -> mitochondria",
], sz=12, sp=Pt(3))
tb(slide, Inches(5.3), Inches(3.5), Inches(2.8), Inches(0.8), "~2 billion years\nendosymbiosis", sz=13, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)

# ===================== SLIDE 4: WHY COMPARE / NAME =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Why Compare? & The Project Name", "A natural experiment spanning 2 billion years of evolution"); footer(slide, sn)
bullets(slide, Inches(0.6), Inches(1.3), Inches(7), Inches(2.8), [
    "If substrate properties are CONSERVED across species:",
    "    -> Strong evolutionary constraints on what makes folding difficult",
    "    -> Structural features requiring chaperonin are fundamental & ancient",
    "",
    "If substrate properties have DIVERGED:",
    "    -> Chaperonin-substrate co-evolution during bacterium -> organelle transition",
    "    -> Mitochondrial import adds new constraints (MTS, unfolding, re-folding in matrix)",
], sz=15, sp=Pt(5))
add_rect(slide, Inches(0.3), Inches(4.4), Inches(12.7), Inches(2.8), VL_ORANGE)
tb(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4), '"Antah Asti Prarambh" = Sanskrit for "The End is the Beginning"', sz=20, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(2.0), [
    "Central question: Does the C-terminus ('end' of translation) represent a new 'beginning' for chaperonin-assisted folding?",
    "Proteins synthesized N-to-C on ribosomes. N-terminal region emerges FIRST, folds FIRST.",
    "C-terminal region emerges LAST, folds LAST — potentially the most vulnerable region.",
    "We test whether N-vs-C structural asymmetry is specific to chaperonin substrates or universal.",
    "Spoiler: It's universal — the 'end' is NOT a new 'beginning' for chaperonins specifically.",
], sz=13, sp=Pt(3))

# ===================== SLIDE 5: THREE GOALS + HYPOTHESES =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Three Scientific Goals & Nine Pre-Registered Hypotheses"); footer(slide, sn)
# Goal boxes
for i, (title, color, bg, items, left_pos) in enumerate([
    ("Goal 1: Domain Architecture", MED_BLUE, VL_BLUE, [
        "Q: Do substrates have distinctive fold topologies?",
        "CATH superfamily distribution comparison",
        "Fisher's exact test, odds ratios + 95% CI",
        "Size-matched proteome background",
    ], Inches(0.3)),
    ("Goal 2: N-vs-C Asymmetry", ACCENT_ORANGE, VL_ORANGE, [
        "Q: Do N-domains have higher contact order?",
        "Three-region decomposition (pre-domain, N, C)",
        "Wilcoxon signed-rank (paired within-protein)",
        "Substrate-specific or universal?",
    ], Inches(4.5)),
    ("Goal 3: MTS Architecture", ACCENT_GREEN, VL_GREEN, [
        "Q: How do targeting signals relate to domains?",
        "HSP60 matrix enrichment analysis",
        "MTS gap to first domain (pre-domain vs overlap)",
        "Binomial test for pre-domain dominance",
    ], Inches(8.7)),
]):
    add_rect(slide, left_pos, Inches(1.3), Inches(4.0), Inches(2.2), bg)
    tb(slide, left_pos + Inches(0.2), Inches(1.35), Inches(3.6), Inches(0.35), title, sz=14, bold=True, color=color)
    bullets(slide, left_pos + Inches(0.2), Inches(1.75), Inches(3.6), Inches(1.6), items, sz=11, sp=Pt(2))

# Hypothesis table
hyp_data = [
    ["Family", "ID", "Hypothesis", "Test", "Effect Size"],
    ["Domain\nArch.", "H1.1", "GroEL enriched for specific CATH superfamilies (TIM barrels, Rossmann)", "Fisher's exact", "Odds ratio + 95% CI"],
    ["", "H1.2", "HSP60 enriched for specific folds vs matrix background", "Fisher's exact", "Odds ratio + 95% CI"],
    ["", "H1.3", "Fold enrichment conserved between GroEL and HSP60", "Chi-squared", "Cramer's V"],
    ["N-vs-C", "H2.1", "N-domains have different contact order than C-regions (paired)", "Wilcoxon signed-rank", "Rank-biserial r"],
    ["", "H2.2", "N-vs-C asymmetry GREATER in substrates than background", "Mann-Whitney U", "Rank-biserial r"],
    ["", "H2.3", "Class III > Class I GroEL asymmetry (class gradient)", "Kruskal-Wallis H", "Eta-squared"],
    ["MTS", "H3.1", "HSP60 substrates enriched for matrix localization", "Fisher's exact", "Odds ratio"],
    ["", "H3.2", "MTS-bearing substrates have distinct first-domain properties", "Mann-Whitney U", "Rank-biserial r"],
    ["", "H3.3", "MTS is predominantly a pre-domain extension (not domain overlap)", "Binomial (H0: p=0.5)", "Observed proportion"],
]
table(slide, Inches(0.2), Inches(3.7), Inches(12.9), Inches(3.5), hyp_data,
      cw=[Inches(0.9), Inches(0.5), Inches(6.0), Inches(2.5), Inches(3.0)])

# ===================== SLIDE 6: SEVEN DATASETS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "The Seven Datasets", "Each dataset serves a specific purpose in the experimental design"); footer(slide, sn)
ds = [
    ["#", "Dataset", "Size", "Source", "Purpose", "Key File"],
    ["1", "E. coli K-12", "4,403", "UniProt UP000000625", "Background for GroEL", "ecoli_k12_proteome.fasta"],
    ["2", "Human", "20,416", "UniProt UP000005640", "Orthology; parent set", "human_proteome.fasta"],
    ["3", "Human Mito", "1,136", "MitoCarta 3.0", "Compartment-matched bg", "human_mito_proteome.tsv"],
    ["4", "GroEL Substrates", "252", "Kerner 2005, Cell (Table S3)", "Primary GroEL set", "groel_substrates_standardized.tsv"],
    ["5", "HSP60 Tier-1", "266", "Morten 2020, Mol Cell", "Primary HSP60 set", "hsp60_tier1_substrates.tsv"],
    ["6", "Homolog Pairs", "69 pairs", "OrthoFinder + RBH", "Cross-species comparison", "groel_hsp60_homologs.tsv"],
    ["7", "Mito Matrix", "525", "MitoCarta 3.0 (matrix)", "Tightest HSP60 bg", "human_matrix_proteome.tsv"],
]
table(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(2.8), ds,
      cw=[Inches(0.4), Inches(1.7), Inches(0.8), Inches(2.8), Inches(3.5), Inches(3.7)])

tb(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.3), "Data Cleaning Details", sz=16, bold=True, color=DARK_BLUE)
bullets(slide, Inches(0.5), Inches(4.6), Inches(6.0), Inches(2.7), [
    "GroEL (scripts/validate_uniprot_accessions.py):",
    "  149/252 accessions demerged (2005 -> 2026 UniProt)",
    "  UniProt REST API: /uniprotkb/{acc}, 3 retries, 5s backoff",
    "  Demerge resolution: filter to taxon 83333 (K-12 MG1655)",
    "  Prefer reviewed Swiss-Prot over TrEMBL",
    "  4 plasmid-specific proteins flagged (not in K-12 ref)",
    "  SCOP fold codes extracted via regex: [a-g]\\.\\d+",
], sz=11, sp=Pt(2))
bullets(slide, Inches(6.8), Inches(4.6), Inches(6.0), Inches(2.7), [
    "HSP60 (scripts/filter_hsp60_interactome.py):",
    "  3 SILAC ratio columns: coIP1, coIP2, coIP3",
    "  NDIC imputation: 2x 95th percentile per column",
    "  Excluded: baits (HSPD1, HSPE1), co-chaperones (TRAP1,",
    "    HSPA9, GRPEL1, DNAJA3), contaminants (IGH/KRT/TUB prefix)",
    "  Tier 1: MitoCarta + median SILAC > 5.0 (n=266)",
    "  MitoCarta v2->v3: 70 reclassifications (52 MIM shifts)",
], sz=11, sp=Pt(2))

# ===================== SLIDE 7: NINE DECISIONS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Nine Critical Methodological Decisions", "Each with explicit scientific rationale documented before analysis"); footer(slide, sn)
dec = [
    ["#", "Decision", "Rationale", "Impact if Wrong"],
    ["1", "CATH/Chainsaw for domain boundaries\n(NOT InterPro/Pfam)", "InterPro = sequence domains (HMM). Study requires\nstructural folding unit boundaries.", "Wrong domain boundaries ->\nincorrect N-vs-C decomposition"],
    ["2", "Contact order + FoldX for stability\n(NOT pLDDT alone)", "pLDDT = AlphaFold confidence, NOT thermodynamic\nstability. CO = kinetics. FoldX DG = thermodynamics.", "Conflate prediction confidence\nwith protein stability"],
    ["3", "OrthoFinder on full proteomes\n(RBH as supplementary)", "RBH misses many-to-many. OrthoFinder captures\nparalogs. 40 -> 69 pairs.", "Miss 29 orthogroup-only\nhomolog pairs"],
    ["4", "SILAC-based HSP60 filtering\n+ MitoCarta confirmation", "Co-IP captures interactions, not function.\nSILAC > 5 + MitoCarta = high confidence.", "Include co-chaperones and\ncontaminants as 'substrates'"],
    ["5", "Compartment + size-matched\ncontrols", "Prevents conflating mito properties\nwith substrate properties.", "Spurious enrichments from\ncompartment/size bias"],
    ["6", "Three-region decomposition\n(pre-domain, N-domain, C-region)", "Transit peptide != N-terminus != first domain.\nConflating produces meaningless results.", "Mix MTS properties with\ndomain stability measurements"],
    ["7", "MitoCarta 3.0 as ground truth\n(predictors only for gaps)", "Experimental evidence > 15-30% error rate\npredictors. 70 reclassifications v2->v3.", "Wrong compartment assignments\nfor 70 proteins"],
    ["8", "Hierarchical BH correction\n(pre-registered hypotheses)", "56 tests need FDR control. Within-family BH +\nacross-family BH (Simes). alpha=0.05.", "Many false positives\nor overly conservative"],
    ["9", "heiniglab/STRIDE binary\n(NOT bioconda 'stride')", "Bioconda stride = genomic StriDe v0.0.1.\nProtein STRIDE = Frishman & Argos.", "Zero output from Chainsaw,\nndom=0 for all proteins"],
]
table(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(6.0), dec,
      cw=[Inches(0.4), Inches(2.5), Inches(5.5), Inches(4.5)])

# ===================== SLIDE 8: TOOLS & SOFTWARE STACK =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Complete Software & Tools Stack", "Every tool with version, source, and purpose"); footer(slide, sn)
sw = [
    ["Tool", "Version", "Source", "Architecture", "Purpose"],
    ["Python", "3.9.16", "Anaconda", "ARM64 (M1 native)", "Core analysis language"],
    ["pandas", "2.2.2", "pip/conda", "-", "Data manipulation & I/O"],
    ["scipy", "1.9.2", "pip/conda", "-", "Statistical tests (Wilcoxon, Mann-Whitney, Fisher, KW)"],
    ["BioPython", "1.78", "pip/conda", "-", "Sequence/structure parsing, CIF->PDB conversion"],
    ["matplotlib + seaborn", "installed", "pip/conda", "-", "Publication figures (300 DPI, PDF+PNG)"],
    ["openpyxl", "installed", "pip/conda", "-", "MitoCarta Excel parsing"],
    ["MMseqs2", "v18.8cc5c", "conda (proteomics env)", "x86_64 (Rosetta 2)", "RBH + all-vs-all sequence search"],
    ["Foldseek", "v10.941cd33", "conda (proteomics env)", "x86_64 (Rosetta 2)", "3Di+AA structural search & clustering"],
    ["Chainsaw", "v3 (2024)", "github/JudeWells", "HPC compiled", "ML domain boundary prediction (CATH-trained)"],
    ["STRIDE", "heiniglab/stride", "Compiled from source", "HPC compiled", "Protein secondary structure assignment"],
    ["mkDSSP", "2.2.1", "conda (base)", "-", "DSSP secondary structure (H/G/I/E/B/T/S/-)"],
    ["FoldX", "5.1 (build 20270131)", "foldxsuite.crg.eu", "Linux x86_64", "Thermodynamic stability (DeltaG kcal/mol)"],
    ["gemmi", "installed", "pip", "-", "CIF file parsing & CIF->PDB conversion"],
    ["Snakemake", "7.32.4", "conda (base)", "-", "Workflow management (Phase 2)"],
    ["AlphaFold DB", "v4 / v6", "EBI FTP", "-", "Predicted 3D structures (CIF format)"],
]
table(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(5.5), sw,
      cw=[Inches(2.0), Inches(1.5), Inches(2.3), Inches(2.0), Inches(5.1)])
tb(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
   "Local: Apple M1 MacBook Air, 8 GB RAM  |  HPC: CSIR-IGIB 'Tejas', SLURM, Lustre, 40 CPUs/380 GB RAM per node, QOS: 5 concurrent jobs",
   sz=11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================== SLIDE 9: PHASE 1 WORKFLOW =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Phase 1: Pilot Workflow — Modules A through I", "1,390 proteins analyzed on local Mac to validate all methods"); footer(slide, sn)
mods = [
    ["Module", "Script(s)", "Input", "Output", "Key Details"],
    ["A: Data Cleaning", "validate_uniprot_accessions.py\nfilter_hsp60_interactome.py\nparse_mitocarta.py", "Raw literature data,\nUniProt, MitoCarta XLS", "Standardized TSVs\nfor all 7 datasets", "149 demerged, SILAC>5,\nMC3.0 binary flags"],
    ["B: Datasets", "module_c_extract_fasta.py", "Proteome FASTAs +\nsubstrate TSVs", "groel_substrates.fasta\nhsp60_tier1.fasta", "API fallback for\nmissing sequences"],
    ["C: Orthology", "module_c_analyze_rbh.py\nrun_orthology.py\nbuild_dataset6_homologs.py", "Substrate FASTAs +\nfull proteome FASTAs", "40 RBH + 62 OrthoFinder\n= 69 merged pairs", "MMseqs2 e-val<1e-5\n%id>25%, cov>50%"],
    ["D: Structures", "download_alphafold_pilot.py\nrun_dssp.py\nvalidate_structure_quality.py", "Accession lists", "1,382 CIFs, DSSP,\nquality tiers", "v6 fallback v4,\npLDDT thresholds"],
    ["E: Domains", "get_cath_domains.py\nrun_chainsaw_e2.py\nanalyze_foldseek.py\ncompute_domain_metrics.py", "CIF structures +\nInterPro API", "1,387/1,390 assigned\n(99.8% coverage)", "Gene3D API 1 req/s\nChainsaw batch 500"],
    ["F: N-vs-C", "module_f_n_vs_c_analysis.py\nmodule_f_extension_chainsaw.py", "Domain boundaries +\nCIF coordinates", "Paired N-vs-C for\nall multi-domain", "RCO: 8A cutoff,\nmin_sep=6 residues"],
    ["G: MTS", "module_g_mts_analysis.py", "UniProt transit peptide +\nMitoCarta + domain boundaries", "MTS-domain gap,\ntargeting classification", "Gap = domain_start -\ntransit_end"],
    ["H: Statistics", "module_h_comparative_stats.py", "All Module E/F/G outputs", "281 tests, 22 significant\n(hierarchical BH)", "alpha=0.05, FDR\ncontrol per family"],
    ["I: Figures", "generate_figures.py", "All results TSVs", "6 figures x (PDF+PNG)\n300 DPI", "Colorblind palette\n(Okabe-Ito)"],
]
table(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(5.8), mods,
      cw=[Inches(1.3), Inches(2.8), Inches(2.2), Inches(2.5), Inches(4.1)])

# ===================== SLIDE 10: MODULE C — ORTHOLOGY PARAMETERS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module C: Orthology Detection — Exact Parameters", "MMseqs2 RBH + OrthoFinder-style connected-component clustering"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.0), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.35), "Step C1: Reciprocal Best Hit (RBH)", sz=16, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.4), [
    "Tool: MMseqs2 v18.8cc5c easy-rbh",
    "Command: mmseqs easy-rbh groel.fasta hsp60.fasta out tmp",
    "Output format: query,target,evalue,bits,alnlen,qcov,tcov,pident,qlen,tlen",
    "E-value: default (no explicit cutoff in easy-rbh mode)",
    "Search space: 252 x 266 = 67,032 comparisons",
    "Result: 40 RBH pairs",
    "Median sequence identity: 35.8%",
    "E-value range: 2.3x10^-156 to 7.8x10^-3",
    "Class III observation: only 8/84 (9.5%) have orthologs",
], sz=11, sp=Pt(2))

add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.0), VL_ORANGE)
tb(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.35), "Step C2-C3: OrthoFinder-Style Orthology", sz=16, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.4), [
    "Tool: MMseqs2 v18.8cc5c (bidirectional all-vs-all)",
    "Search parameters:",
    "    E-value cutoff: 1e-5 (0.00001)",
    "    Min percent identity: 25.0%",
    "    Min coverage: 50% (query OR target)",
    "    Sensitivity: 7.5",
    "    Max sequences: 300",
    "    Threads: 4",
    "Search space: 4,403 x 20,416 = ~89.9 million",
    "Clustering: Union-Find connected components",
    "Edge = reciprocal hit (found in BOTH directions)",
    "Result: 422 orthogroups, 34 shared substrate groups",
], sz=11, sp=Pt(2))

add_rect(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.7), VL_GREEN)
tb(slide, Inches(0.5), Inches(4.6), Inches(12), Inches(0.35), "Step C4: Dataset 6 Merging (build_dataset6_homologs.py)", sz=16, bold=True, color=ACCENT_GREEN)
bullets(slide, Inches(0.5), Inches(5.0), Inches(5.8), Inches(2.0), [
    "Merge logic:",
    "  1. Expand orthogroup many-to-many into all pairs",
    "  2. Match RBH pairs against orthogroup pairs",
    "  3. Label evidence: 'both', 'rbh_only', 'orthogroup_only'",
    "  4. Sort: both > rbh > orthogroup, then by e-value",
    "Result: 69 unique pairs",
    "    33 found by both methods (highest confidence)",
    "    7 RBH-only pairs",
    "    29 orthogroup-only pairs",
], sz=11, sp=Pt(2))
bullets(slide, Inches(6.8), Inches(5.0), Inches(5.8), Inches(2.0), [
    "Validation:",
    "  79.7% (55/69) share same top CATH superfamily",
    "  Fold conservation confirms biological relevance",
    "",
    "Why both methods?",
    "  RBH: simple, 1-to-1, but misses paralogs",
    "  OrthoFinder: captures gene families but noisier",
    "  Merging gives best of both worlds",
    "  Evidence column allows sensitivity analyses",
], sz=11, sp=Pt(2))

# ===================== SLIDE 11: MODULE D — ALPHAFOLD + DSSP =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module D: Structure Acquisition & Quality Control", "AlphaFold download, DSSP secondary structure, quality validation"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.1), Inches(3.0), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(3.7), Inches(0.3), "D1-D2: AlphaFold Download", sz=14, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.75), Inches(3.7), Inches(2.4), [
    "URL: alphafold.ebi.ac.uk/files/AF-{acc}-F1-model_v6.cif",
    "Fallback: v4 if v6 unavailable",
    "Batch size: 50 proteins, 0.5s delay between batches",
    "Per-request timeout: 30 seconds",
    "Streaming chunk: 65,536 bytes",
    "Max retries: 3 (exponential backoff)",
    "Pilot: 1,382/1,390 structures (99.4%)",
    "8 missing: P07203, P30042, P36969, Q16881,",
    "    Q5THJ4, Q86UA3, Q9BVL4, Q9NNW7",
    "CIF parsing: CA atoms only, B-factor = pLDDT",
], sz=10, sp=Pt(1))

add_rect(slide, Inches(4.6), Inches(1.3), Inches(4.1), Inches(3.0), VL_ORANGE)
tb(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(0.3), "D3: DSSP Secondary Structure", sz=14, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(4.8), Inches(1.75), Inches(3.7), Inches(2.4), [
    "Tool: mkdssp v2.2.1",
    "Timeout: 30 seconds per structure",
    "DSSP code mapping:",
    "    Helix:  H (alpha), G (3_10), I (pi)",
    "    Strand: E (beta-strand), B (beta-bridge)",
    "    Coil:   T (turn), S (bend), - (coil), ' ' (none)",
    "Result: 1,382/1,382 processed",
    "Mean composition:",
    "    Helix: 43.5%, Strand: 14.2%, Coil: 42.2%",
    "Per-residue TSV for regional analysis",
], sz=10, sp=Pt(1))

add_rect(slide, Inches(8.9), Inches(1.3), Inches(4.1), Inches(3.0), VL_GREEN)
tb(slide, Inches(9.1), Inches(1.4), Inches(3.7), Inches(0.3), "D4: Quality Validation", sz=14, bold=True, color=ACCENT_GREEN)
bullets(slide, Inches(9.1), Inches(1.75), Inches(3.7), Inches(2.4), [
    "Quality tiers (by mean pLDDT):",
    "    Very High: >= 90 (backbone + sidechains)",
    "    High:      80-90 (reliable backbone)",
    "    Moderate:  70-80 (use with caution)",
    "    Low:       50-70 (fold-level only)",
    "    Very Low:  < 50 (unreliable)",
    "Quality flags:",
    "    flag_very_low_plddt: mean < 50",
    "    flag_majority_unreliable: mean<70 & frac>70<0.5",
    "    flag_few_usable: frac_plddt_gt70 < 0.3",
    "Result: 77.4% high/very high; 63 flagged (4.6%)",
], sz=10, sp=Pt(1))

# Phase 2 download
add_rect(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.7), RGBColor(0xF5, 0xF5, 0xF5))
tb(slide, Inches(0.5), Inches(4.6), Inches(12), Inches(0.3), "Phase 2: Full-Scale Download (download_alphafold_full.py)", sz=14, bold=True, color=DARK_BLUE)
bullets(slide, Inches(0.5), Inches(5.0), Inches(6.0), Inches(2.0), [
    "Bulk FTP download:",
    "  E. coli: ftp.ebi.ac.uk/.../UP000000625_83333_ECOLI_v6.tar (~2 GB)",
    "  Human:   ftp.ebi.ac.uk/.../UP000005640_9606_HUMAN_v6.tar (~20 GB)",
    "  Individual API fallback for missing accessions",
    "  Download workers: 8 threads (config.yaml)",
    "  Checkpoint interval: 500 proteins (JSON-based resume)",
], sz=11, sp=Pt(2))
bullets(slide, Inches(6.8), Inches(5.0), Inches(6.0), Inches(2.0), [
    "Results:",
    "  25,007 structures downloaded total",
    "  E. coli: 4,371 CIF files (99.3% coverage)",
    "  Human: 20,636 CIF files (~100% coverage)",
    "  Total size: ~22 GB on HPC Lustre filesystem",
    "  Bug fix: human CIFs required gunzip (05a_decompress.sh)",
], sz=11, sp=Pt(2))

# ===================== SLIDE 12: MODULE E — DOMAINS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module E: Structural Domain Assignment — Parameters", "CATH Gene3D + Chainsaw ML + Foldseek Structural Clustering"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.1), Inches(3.2), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(3.7), Inches(0.3), "E1: CATH via Gene3D (InterPro API)", sz=14, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.75), Inches(3.7), Inches(2.6), [
    "API: ebi.ac.uk/interpro/api/entry/cathgene3d/",
    "    protein/uniprot/{acc}?format=json",
    "Rate limit: 1.0 req/sec (MIN_DELAY=1.0s)",
    "HTTP timeout: 30 seconds",
    "Max retries: 3, backoff: 5s (exponential)",
    "429/408 status -> automatic retry with backoff",
    "Checkpoint: every 50 proteins",
    "CATH code format: G3DSA:X.Y.Z.W",
    "  X=Class, Y=Arch, Z=Topology, W=Homologous SF",
    "Class mapping: 1=alpha, 2=beta, 3=ab, 4=few-SS",
    "Handles discontinuous domains (multi-segment)",
    "Result: 1,151 proteins, 2,141 domains total",
], sz=10, sp=Pt(1))

add_rect(slide, Inches(4.6), Inches(1.3), Inches(4.1), Inches(3.2), VL_ORANGE)
tb(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(0.3), "E2: Chainsaw ML Predictions", sz=14, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(4.8), Inches(1.75), Inches(3.7), Inches(2.6), [
    "Tool: Chainsaw v3 (Wells et al., 2024)",
    "Trained on: CATH domain structures",
    "Input: AlphaFold CIF + STRIDE secondary structure",
    "STRIDE: heiniglab/stride (compiled from source!)",
    "    NOT bioconda stride (genomic variant caller)",
    "Batch processing: 500 proteins per batch",
    "    Symlink CIFs to temp directory on Lustre",
    "Command: python3 get_predictions.py",
    "    --structure_directory TMPDIR",
    "    --output BATCH.tsv --use_first_chain",
    "Chopping format: '12-42_308-396,56-302'",
    "    Comma = domains, underscore = segments",
    "Phase 2: 25,007 proteins, 93.6% assigned",
], sz=10, sp=Pt(1))

add_rect(slide, Inches(8.9), Inches(1.3), Inches(4.1), Inches(3.2), VL_GREEN)
tb(slide, Inches(9.1), Inches(1.4), Inches(3.7), Inches(0.3), "E4: Foldseek Structural Clustering", sz=14, bold=True, color=ACCENT_GREEN)
bullets(slide, Inches(9.1), Inches(1.75), Inches(3.7), Inches(2.6), [
    "Tool: Foldseek v10.941cd33",
    "Parameters (config.yaml):",
    "    Min sequence identity: 0.3 (30%)",
    "    Coverage: 0.5 (50%, bidirectional)",
    "    Coverage mode: 0 (bidirectional)",
    "    Cluster mode: 0 (set-cover cascaded)",
    "    E-value (search): 0.001",
    "    E-value (cluster): 0.01",
    "    Sensitivity: 7.5",
    "    Alignment type: 2 (3Di+AA hybrid)",
    "    Backtrace: enabled (-a flag)",
    "    Threads: 16, Memory: 64 GB",
    "Pilot: 1,155 clusters, 24 shared (GroEL+HSP60)",
    "Phase 2: 16,193 clusters, 27,063 proteins",
], sz=10, sp=Pt(1))

# Coverage summary
cov = [
    ["Scale", "CATH (Gene3D)", "Chainsaw ML", "Unassigned", "Total Coverage"],
    ["Pilot (1,390)", "1,151 (82.8%)", "236 (17.0%)", "3 (0.2%)", "1,387 (99.8%)"],
    ["Full (25,007)", "1,390 (CATH-only)", "23,868 (full)", "-", "25,258 unified records"],
]
table(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.9), cov)

# Per-domain metrics
bullets(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(1.3), [
    "E3: Per-Domain Structural Metrics (compute_domain_structural_metrics.py): For each domain -> frac_helix, frac_strand, frac_coil, mean_pLDDT, min_pLDDT, frac_plddt_gt70, frac_plddt_gt90",
    "CIF cache: max 500 proteins in memory, evict oldest 100 when full. Mean domain pLDDT: 92.1 (higher than whole-protein because domains exclude disordered regions)",
    "E5: Domain Distribution (domain_distribution_summary.py): Alpha-beta dominates all datasets (60-71%). Top 10 superfamilies per dataset. Domain count: 0-domain/single/multi classification",
], sz=11, sp=Pt(3))

# ===================== SLIDE 13: MODULE F — CONTACT ORDER =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module F: N-vs-C Stability Analysis — Contact Order Parameters", "Plaxco-Simons-Baker (1998) definition with exact implementation details"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.5), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.3), "Contact Order Calculation (Exact Parameters)", sz=16, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.9), [
    "Definition: Plaxco, Simons & Baker, J Mol Biol (1998)",
    "",
    "Absolute CO (ACO):",
    "    ACO = (1/N_contacts) x SUM |i - j|",
    "    for all pairs (i,j) where:",
    "        CA-CA distance < 8.0 Angstroms",
    "        |i - j| >= 6 residues (min sequence separation)",
    "",
    "Relative CO (RCO):",
    "    RCO = ACO / N_residues",
    "",
    "Parameters:",
    "    Distance cutoff: 8.0 A (CA-CA Euclidean)",
    "    Min sequence separation: 6 residues",
    "    Atom type: CA (alpha-carbon) only",
    "    Coordinates: from AlphaFold CIF B-factor field",
    "    Large protein chunking: 1000 residues (n>3000)",
], sz=11, sp=Pt(1))

add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.5), VL_ORANGE)
tb(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.3), "Three-Region Decomposition", sz=16, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.9), [
    "For each multi-domain protein (>= 2 CATH/Chainsaw domains):",
    "",
    "Region 1 — Pre-domain tail:",
    "    Residues 1 to (first_domain_start - 1)",
    "    May include: transit peptide, signal peptide, unstructured N-term",
    "",
    "Region 2 — N-domain (first structural domain):",
    "    Residues first_domain_start to first_domain_end",
    "    Sorted by domain_start position from CATH/Chainsaw",
    "",
    "Region 3 — C-region (everything after first domain):",
    "    Residues (first_domain_end + 1) to protein_length",
    "    Contains all subsequent domains + linkers",
    "",
    "Minimum region size: >= 5 residues (skip smaller)",
    "Single-domain: N-half vs C-half split at midpoint",
], sz=11, sp=Pt(1))

# Metrics computed
add_rect(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(2.2), RGBColor(0xF5,0xF5,0xF5))
tb(slide, Inches(0.5), Inches(5.1), Inches(12), Inches(0.3), "All Metrics Computed Per Region", sz=14, bold=True, color=DARK_BLUE)
met = [
    ["Category", "Metrics", "Source"],
    ["Sequence", "length, net_charge (K+R-D-E), frac_hydrophobic (A,V,I,L,F,W,M), frac_charged (K,R,D,E), frac_polar, frac_aromatic, frac_small (G,A,S)", "Amino acid sequence"],
    ["Structure", "frac_helix (H,G,I), frac_strand (E,B), frac_coil (all others), mean_pLDDT, frac_plddt_gt70, frac_plddt_gt90", "DSSP + CIF B-factors"],
    ["Folding", "absolute_contact_order, relative_contact_order (RCO), n_contacts", "CIF CA coordinates"],
    ["Stability", "FoldX total_energy (DeltaG, kcal/mol) — IN PROGRESS", "FoldX 5.1 RepairPDB + Stability"],
]
table(slide, Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.6), met, cw=[Inches(1.2), Inches(8.0), Inches(3.5)])

# ===================== SLIDE 14: MODULE G — MTS ANALYSIS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module G: Mitochondrial Targeting Signal Analysis", "UniProt feature extraction, MitoCarta integration, gap calculation"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.1), Inches(2.8), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(3.7), Inches(0.3), "G1-G2: UniProt Feature Extraction", sz=14, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.75), Inches(3.7), Inches(2.2), [
    "API: rest.uniprot.org/uniprotkb/stream",
    "Format: TSV",
    "Fields: accession, ft_transit, ft_signal,",
    "    cc_subcellular_location",
    "Batch size: 100 proteins per request",
    "Timeout: 120 seconds",
    "Max retries: 3 (exponential backoff)",
    "Rate limit: 1.0s inter-batch delay",
    "Transit peptide regex: TRANSIT (\\d+)\\.\\.(\\d+)",
    "Signal peptide regex: SIGNAL (\\d+)\\.\\.(\\d+)",
    "Cached in: uniprot_transit_signal_cache.tsv",
], sz=10, sp=Pt(1))

add_rect(slide, Inches(4.6), Inches(1.3), Inches(4.1), Inches(2.8), VL_ORANGE)
tb(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(0.3), "G3-G4: Targeting Classification", sz=14, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(4.8), Inches(1.75), Inches(3.7), Inches(2.2), [
    "MitoCarta 3.0 parsing (parse_mitocarta.py):",
    "    Sheet: 'A Human MitoCarta3.0'",
    "    Binary flags (case-insensitive string match):",
    "        is_matrix: contains 'Matrix'",
    "        is_im: contains 'MIM' or 'Membrane'",
    "        is_ims: contains 'IMS'",
    "        is_om: contains 'MOM'",
    "",
    "Classification hierarchy:",
    "    1. High-confidence matrix: is_matrix AND has_tp",
    "    2. Non-canonical matrix: is_matrix AND !has_tp",
    "    3. Inner membrane / IMS / Outer membrane",
    "    4. Mito with MTS / Mito other",
    "    5. Non-mitochondrial (default)",
], sz=10, sp=Pt(1))

add_rect(slide, Inches(8.9), Inches(1.3), Inches(4.1), Inches(2.8), VL_GREEN)
tb(slide, Inches(9.1), Inches(1.4), Inches(3.7), Inches(0.3), "G5: MTS-Domain Gap Calculation", sz=14, bold=True, color=ACCENT_GREEN)
bullets(slide, Inches(9.1), Inches(1.75), Inches(3.7), Inches(2.2), [
    "For proteins with transit peptide AND CATH domain:",
    "",
    "gap_length = first_domain_start - transit_end",
    "",
    "Classification:",
    "    gap >= 0: MTS is PRE-DOMAIN (separate extension)",
    "        MTS cleaved BEFORE domain starts",
    "        Domain emerges unfolded in matrix",
    "    gap < 0: MTS OVERLAPS domain",
    "        Cleavage disrupts domain integrity",
    "",
    "Result: 84.4% pre-domain (n=368/436)",
    "Median gap: 18 residues",
    "Mean gap: 37.4 residues (range: 0-579)",
], sz=10, sp=Pt(1))

# MitoCarta data cleaning
add_rect(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.9), RGBColor(0xF5,0xF5,0xF5))
tb(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.3), "MitoCarta 3.0 Data Details", sz=14, bold=True, color=DARK_BLUE)
bullets(slide, Inches(0.5), Inches(4.8), Inches(6.0), Inches(2.2), [
    "Source: Broad Institute, Human.MitoCarta3.0.xls",
    "URL: personal.broadinstitute.org/scalvo/MitoCarta3.0/",
    "Total proteins: 1,136 (full mito), 525 (matrix subset)",
    "v2.0 -> v3.0 changes: 70 reclassifications",
    "    52 respiratory chain subunits: Matrix -> MIM",
    "    18 non-Matrix -> Matrix",
    "    Impact: significantly changes background composition",
], sz=11, sp=Pt(2))
bullets(slide, Inches(6.8), Inches(4.8), Inches(6.0), Inches(2.2), [
    "HSP60 cross-reference:",
    "    Lookup by UniProt ID first, then gene symbol (UPPER)",
    "    Identify MC3 gains/losses/localization changes",
    "    HSP60 Tier-1: 266 proteins",
    "    46.6% high-confidence matrix",
    "    21.1% non-canonical matrix import (no MTS)",
    "    Remaining: other mito compartments or non-mito",
], sz=11, sp=Pt(2))

# ===================== SLIDE 15: MODULE H — STATISTICS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module H: Statistical Testing Framework — Exact Methods", "Hierarchical Benjamini-Hochberg correction across 3 hypothesis families"); footer(slide, sn)

# Correction method
add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.3), "Multiple Testing Correction", sz=15, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.75), Inches(5.8), Inches(2.2), [
    "Framework: Hierarchical Benjamini-Hochberg (BH)",
    "Overall alpha: 0.05 (FDR control)",
    "",
    "Level 1: Within each family -> BH correction",
    "    Family 1 (Domain arch.): 24 tests -> BH",
    "    Family 2 (N-vs-C): 30 tests -> BH",
    "    Family 3 (MTS targeting): 2 tests -> BH",
    "",
    "Level 2: Across families -> Simes method",
    "    Simes p = min(k * p(i) / i) for ordered p-values",
    "    BH correction on 3 family-level Simes p-values",
    "",
    "Significant = pass BOTH levels",
], sz=11, sp=Pt(1))

# Test details
add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8), VL_ORANGE)
tb(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.3), "Effect Size Formulas", sz=15, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(7.0), Inches(1.75), Inches(5.8), Inches(2.2), [
    "Wilcoxon signed-rank (paired):",
    "    r = (R+ - R-) / (n(n+1)/2)   [rank-biserial]",
    "",
    "Mann-Whitney U (between-group):",
    "    r = 1 - (2U) / (n1 * n2)     [rank-biserial]",
    "",
    "Kruskal-Wallis (3+ groups):",
    "    eta^2 = (H - k + 1) / (N - k)",
    "",
    "Fisher's exact (enrichment):",
    "    OR = (a*d) / (b*c)  [Woolf's 95% CI]",
    "",
    "Chi-squared (distribution):",
    "    Cramer's V = sqrt(chi2 / (n * (min(r,c)-1)))",
], sz=11, sp=Pt(1))

# Test list
tests = [
    ["Family", "Test ID", "Comparison", "Method", "Metrics Tested"],
    ["1: Domain", "H1.1-H1.2", "Substrate enrichment per CATH SF", "Fisher's exact (2x2)", "Top 10 SFs per dataset"],
    ["", "H1.3", "CATH class distribution", "Chi-squared", "4-class proportions"],
    ["2: N-vs-C", "H2.1 (x4 datasets)", "N-domain vs C-region (paired)", "Wilcoxon signed-rank", "RCO, pLDDT, helix, strand, plddt>70, hydrophobic"],
    ["", "H2.2 (x4)", "Substrate vs background (N-C diff)", "Mann-Whitney U", "RCO difference, pLDDT difference"],
    ["", "H2.3 (x2)", "GroEL Class I/II/III", "Kruskal-Wallis H", "RCO difference, pLDDT difference"],
    ["3: MTS", "H3.1", "HSP60 matrix enrichment", "Fisher's exact", "Matrix vs non-matrix in mito"],
    ["", "H3.3", "MTS pre-domain dominance", "Binomial (H0: p=0.5)", "Proportion pre-domain vs overlap"],
]
table(slide, Inches(0.2), Inches(4.3), Inches(12.9), Inches(3.0), tests,
      cw=[Inches(1.0), Inches(1.5), Inches(3.5), Inches(2.5), Inches(4.4)])

# ===================== SLIDE 16: PHASE 2 HPC PIPELINE =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Phase 2: HPC Pipeline — SLURM Resource Allocation", "CSIR-IGIB HPC 'Tejas': 19 SLURM job scripts with dependency chains"); footer(slide, sn)

slurm = [
    ["Step", "Job Name", "Script", "CPUs", "Memory", "Wall Time", "Key Command / Purpose"],
    ["0", "aap_setup", "00_setup.sh", "1", "2 GB", "1h", "Validate env, create dirs, check 50GB disk"],
    ["1", "aap_download", "01_download.sh", "1", "4 GB", "12h", "download_alphafold_full.py --workers 8"],
    ["2", "aap_fs_createdb", "02_foldseek.sh", "4", "16 GB", "2h", "run_foldseek_full.py --step createdb"],
    ["3", "aap_fs_search", "03_foldseek.sh", "16", "64 GB", "24h", "run_foldseek_full.py --step search (CRITICAL)"],
    ["4", "aap_fs_cluster", "04_foldseek.sh", "16", "64 GB", "8h", "cluster + export + analyze-only"],
    ["5", "aap_chainsaw", "05_chainsaw.sh", "4", "16 GB", "72h", "Batch 500, symlinks on Lustre (CRITICAL)"],
    ["5a", "-", "05a_decompress.sh", "-", "-", "-", "gunzip human CIF.gz files (bug fix)"],
    ["6a", "aap_foldx_gen", "06_foldx_gen.sh", "1", "2 GB", "30m", "run_foldx.py --generate-slurm (501 tasks)"],
    ["6b", "(array)", "foldx_array.slurm", "1", "4 GB", "6h/task", "MANUAL sbatch, 501 tasks x 50 proteins"],
    ["7", "aap_foldx_col", "07_foldx_col.sh", "1", "4 GB", "1h", "run_foldx.py --collect -> merged TSV"],
    ["8", "aap_mod_e", "08_module_e.sh", "4", "16 GB", "4h", "Unified CATH + Chainsaw domains"],
    ["9", "aap_mod_f", "09_module_f.sh", "4", "16 GB", "8h", "N-vs-C stability + contact order"],
    ["10", "aap_mod_h", "10_module_h.sh", "4", "16 GB", "4h", "56 tests, hierarchical BH correction"],
    ["11", "aap_mod_i", "11_module_i.sh", "2", "8 GB", "2h", "6 figures, MPLBACKEND=Agg"],
]
table(slide, Inches(0.1), Inches(1.2), Inches(13.1), Inches(5.5), slurm,
      cw=[Inches(0.5), Inches(1.3), Inches(1.5), Inches(0.6), Inches(0.8), Inches(0.9), Inches(7.5)])

tb(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
   "ALL scripts include: export LD_LIBRARY_PATH=$CONDA_PREFIX/lib:$LD_LIBRARY_PATH (gcc-8.4/numpy fix) + python3 -u (unbuffered for real-time logs)",
   sz=11, bold=True, color=ACCENT_RED)

# ===================== SLIDE 17: PIPELINE DEPENDENCY + FOLDX =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Pipeline Dependency Graph & FoldX Parameters", "Parallel branches + sequential analysis chain; FoldX thermodynamic stability"); footer(slide, sn)

# Dependency text diagram
add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.5), Inches(3.5), VL_BLUE)
tb(slide, Inches(0.5), Inches(1.4), Inches(6.1), Inches(0.3), "Dependency Graph", sz=15, bold=True, color=MED_BLUE)
bullets(slide, Inches(0.5), Inches(1.8), Inches(6.1), Inches(2.9), [
    "00_setup -> 01_download (afterok)",
    "01_download ->",
    "    Branch A: 02_createdb -> 03_search -> 04_cluster",
    "    Branch B: 05_chainsaw (parallel with A)",
    "    Branch C: 06_foldx_generate (parallel with A,B)",
    "        -> [MANUAL] sbatch foldx_array.slurm",
    "        -> 07_foldx_collect",
    "",
    "08_module_e (afterok: 04_cluster AND 05_chainsaw)",
    "",
    "Analysis chain (submit_analysis.sh):",
    "    09_module_f (afterok: foldx_collect or 'done')",
    "        -> 10_module_h (afterok: 09)",
    "            -> 11_module_i (afterok: 10)",
    "",
    "QOS limit: max 5 concurrent jobs (common QOS)",
], sz=11, sp=Pt(1))

# FoldX parameters
add_rect(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(3.5), VL_ORANGE)
tb(slide, Inches(7.2), Inches(1.4), Inches(5.6), Inches(0.3), "FoldX 5.1 Thermodynamic Stability", sz=15, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(7.2), Inches(1.8), Inches(5.6), Inches(2.9), [
    "Binary: /lustre/vishal.bharti/software/foldx5/foldx",
    "    (symlink to foldx_20270131)",
    "License: Academic, expires 2026-12-31",
    "Rotabase: NOT required (compiled into FoldX 5.1)",
    "",
    "Thermodynamic parameters:",
    "    Temperature: 298.15 K (25C)",
    "    pH: 7.0",
    "    Ionic strength: 0.05 M",
    "",
    "Pipeline per protein:",
    "    1. CIF -> PDB conversion (gemmi or BioPython)",
    "    2. FoldX RepairPDB (optimize rotamers, fix clashes)",
    "    3. FoldX Stability (compute DeltaG at above params)",
    "    4. Save as JSON: per_protein/{ACC}.json",
    "",
    "Output: total_energy (DeltaG, kcal/mol)",
    "    Component energies: backbone_hbond, sidechain_hbond,",
    "    vdw_clashes, electrostatics, solvation_polar/hydrophobic,",
    "    entropy_mainchain/sidechain (may be null in 5.1)",
    "",
    "Array: 501 tasks x 50 proteins, timeout 300s/protein",
    "Status: ~42% complete, est. April 1-2, 2026",
], sz=10, sp=Pt(0))

# Biological rationale
add_rect(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(2.2), RGBColor(0xF5,0xF5,0xF5))
tb(slide, Inches(0.5), Inches(5.1), Inches(12), Inches(0.3), "Why Contact Order AND FoldX? (Complementary Metrics)", sz=14, bold=True, color=DARK_BLUE)
bullets(slide, Inches(0.5), Inches(5.5), Inches(6.0), Inches(1.5), [
    "Contact Order = KINETICS (folding speed)",
    "    High CO -> folds SLOWLY (many long-range contacts)",
    "    Correlates with experimental kf (r ~ -0.75)",
    "    Structural property from 3D coordinates",
], sz=12, sp=Pt(2))
bullets(slide, Inches(6.8), Inches(5.5), Inches(6.0), Inches(1.5), [
    "FoldX DeltaG = THERMODYNAMICS (fold stability)",
    "    More negative DG -> MORE STABLE folded state",
    "    Accounts for: vdW, solvation, H-bonds, electrostatics",
    "    High CO + negative DG = ideal chaperonin substrate",
], sz=12, sp=Pt(2))

# ===================== SLIDE 18: FIGURE GENERATION PARAMETERS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Module I: Publication Figure Generation — Parameters", "6 figures, colorblind-friendly, 300 DPI, PDF+PNG"); footer(slide, sn)

bullets(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.5), [
    "Scripts: generate_figures.py (Phase 1)",
    "         module_i_polished.py (Phase 2, local polish)",
    "",
    "Matplotlib settings:",
    "    Default font size: 12pt",
    "    Legend font: 10pt",
    "    DPI: 300 (both screen and savefig)",
    "    PDF/PS font type: 42 (TrueType subset)",
    "    Backend on HPC: Agg (non-interactive)",
    "",
    "Color palette: Colorblind-friendly (Okabe-Ito / seaborn Set2)",
    "    GroEL: blue | HSP60: orange | Matrix: green | Mito: red",
    "    N-domain: blue | C-region: orange",
], sz=12, sp=Pt(2))

fig_table = [
    ["Figure", "Size", "Content", "Plot Types"],
    ["Fig 1: Domain Architecture", "16x5.5 in", "CATH class dist., top SFs, domain counts", "Stacked bar, horizontal bar, grouped bar"],
    ["Fig 2: N-vs-C Stability", "18x6 in", "RCO violin, pLDDT violin, N-C difference heatmap", "Split violin (cut=0, inner=quart), heatmap (diverging)"],
    ["Fig 3: GroEL Class", "12x5.5 in", "RCO diff by class, pLDDT diff by class", "Boxplot + strip (jitter), KW p-value annotation"],
    ["Fig 4: MTS Targeting", "~14x5 in", "Localization breakdown, gap histogram, scatter", "Horizontal bar, histogram, scatter (alpha=0.15)"],
    ["Fig 5: Orthology", "~14x5 in", "Orthogroup categories, RCO conservation", "Bar chart, scatter + regression (r, p annotated)"],
    ["Fig 6: Summary", "~16x5 in", "Key findings overview, MTS pie chart", "Bar + error bars, stacked bar, pie chart"],
]
table(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(2.5), fig_table,
      cw=[Inches(2.5), Inches(1.3), Inches(5.0), Inches(3.9)])

bullets(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5), [
    "All figures annotated with: real p-values (not placeholder), sample sizes (n=X), effect sizes. Both PDF (vector for journals) and PNG (raster for presentations) at 300 DPI.",
], sz=12, sp=Pt(2))

# ===================== SLIDE 19: RESULTS — DOMAIN ARCHITECTURE =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Results: Goal 1 — Domain Architecture", "Chaperonin substrates enriched for specific complex fold topologies"); footer(slide, sn)
fig1 = os.path.join(FIG_DIR, "fig1_domain_architecture.png")
if os.path.exists(fig1): slide.shapes.add_picture(fig1, Inches(0.3), Inches(1.2), Inches(12.7), Inches(3.5))
tb(slide, Inches(0.5), Inches(4.9), Inches(12), Inches(0.3), "Key Enrichments (Phase 2, Fisher's exact, size-matched controls, BH corrected)", sz=14, bold=True, color=DARK_BLUE)
enr = [
    ["Dataset", "CATH Code", "Name", "Odds Ratio", "95% CI", "p (BH)", "Interpretation"],
    ["GroEL", "3.20.20.70", "TIM barrel / Aldolase I", "8.4", "[3.83, 18.30]", "1.9x10^-8", "8-stranded beta/alpha barrel"],
    ["GroEL", "1.10.10.10", "Winged helix-like", "50.9", "[6.3, 412]", "3.5x10^-9", "Complex helix-turn-helix"],
    ["GroEL", "3.30.420.40", "Nucleotidyltransferase", "6.0", "-", "5.8x10^-3", "Enzyme fold"],
    ["GroEL", "3.40.640.10", "Muconolactone isomerase", "2.6", "-", "1.5x10^-2", "Metabolic enzyme"],
    ["HSP60", "3.30.830.10", "Rossmann-like", "5.4", "-", "2.5x10^-4", "Matrix enzyme fold"],
    ["HSP60", "3.90.226.10", "-", "4.8", "-", "4.6x10^-4", "Mitochondrial fold"],
    ["HSP60", "3.40.50.620", "-", "3.3", "-", "1.3x10^-2", "Alpha-beta fold"],
    ["HSP60", "2.40.30.10", "-", "3.6", "-", "1.2x10^-2", "Beta sandwich"],
]
table(slide, Inches(0.2), Inches(5.3), Inches(12.9), Inches(2.0), enr,
      cw=[Inches(0.8), Inches(1.3), Inches(2.5), Inches(1.0), Inches(1.3), Inches(1.2), Inches(4.8)])

# ===================== SLIDE 20: WHY TIM BARRELS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Why TIM Barrels Need Chaperonins", "The topology argument: 8-fold symmetry requires all elements present before folding"); footer(slide, sn)
bullets(slide, Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.5), [
    "TIM barrel = CATH 3.20.20 (Topology level)",
    "    8 alternating beta-strands and alpha-helices",
    "    Forming a closed barrel with 8-fold symmetry",
    "",
    "Why it CAN'T fold co-translationally:",
    "    1. The barrel requires ALL 8 beta-strands to close",
    "    2. During translation, only part of the barrel is synthesized",
    "    3. Partial barrel = exposed hydrophobic core = AGGREGATION",
    "    4. Protein MUST be fully synthesized before barrel can form",
    "",
    "This creates a vulnerable window:",
    "    Full-length unfolded TIM barrel exposes hydrophobic residues",
    "    Without GroEL, it aggregates irreversibly (Class III behavior)",
    "    GroEL barrel (~85 A cavity) provides protected folding space",
    "",
    "GroEL enrichment: OR = 8.4 (Phase 2, p = 1.9x10^-8)",
    "    TIM barrel substrates are 8.4x more common in GroEL",
    "    than expected from the E. coli cytoplasmic proteome",
    "",
    "Multi-domain enrichment: NOT significant (OR=1.13, p=0.35)",
    "    -> It's about FOLD TOPOLOGY, not domain count",
], sz=13, sp=Pt(3))
# Chi-squared result
bullets(slide, Inches(7.5), Inches(1.3), Inches(5.5), Inches(5.5), [
    "CATH class distribution (Chi-squared):",
    "    GroEL: chi2=16.79, p=2.1x10^-3, V=0.089",
    "        -> Significant class effect",
    "    HSP60: chi2=9.23, p=0.055",
    "        -> Not significant (borderline)",
    "",
    "Phase 1 vs Phase 2 consistency:",
    "",
    "    TIM barrel OR:",
    "        Phase 1 (pilot): 9.16",
    "        Phase 2 (full):  8.4",
    "        -> ROBUST (confirmed at full scale)",
    "",
    "    Winged helix OR:",
    "        Phase 1: 42.8",
    "        Phase 2: 50.9",
    "        -> Even STRONGER with larger background",
    "",
    "Conclusion: Chaperonin substrate identity is",
    "determined by fold topology, not general",
    "structural complexity or domain count.",
], sz=13, sp=Pt(3))

# ===================== SLIDE 21: N-vs-C RESULTS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Results: Goal 2 — N-vs-C Contact Order Asymmetry", "N-domains have HIGHER contact order — but it's UNIVERSAL, not substrate-specific"); footer(slide, sn)
fig2 = os.path.join(FIG_DIR, "fig2_n_vs_c_stability.png")
if os.path.exists(fig2): slide.shapes.add_picture(fig2, Inches(0.2), Inches(1.2), Inches(12.9), Inches(3.3))
rco = [
    ["Dataset", "n pairs", "Median N-C diff", "Effect r", "p-value", "Significant?", "Interpretation"],
    ["GroEL substrates", "124", "0.043", "0.41", "8.9x10^-5", "YES", "Medium effect, N>C"],
    ["HSP60 substrates", "131", "0.059", "0.46", "5.3x10^-6", "YES", "Medium-large effect"],
    ["Matrix background", "251", "0.069", "0.43", "2.4x10^-9", "YES", "Same effect in bg!"],
    ["Mito background", "425", "0.064", "0.48", "7.1x10^-18", "YES (STRONGEST)", "Largest N and effect"],
]
table(slide, Inches(0.2), Inches(4.7), Inches(12.9), Inches(1.8), rco,
      cw=[Inches(1.8), Inches(0.8), Inches(1.5), Inches(0.8), Inches(1.2), Inches(1.5), Inches(5.3)])
add_rect(slide, Inches(0.3), Inches(6.7), Inches(12.7), Inches(0.5), VL_RED)
tb(slide, Inches(0.5), Inches(6.75), Inches(12), Inches(0.4),
   "KEY INSIGHT: The mito BACKGROUND shows the strongest effect (r=0.48, p=7.1x10^-18) — this is NOT a chaperonin-specific feature!", sz=14, bold=True, color=ACCENT_RED)

# ===================== SLIDE 22: CRUCIAL NEGATIVE RESULTS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "The Crucial Negative Results — H2.2 & H2.3 REJECTED", "What we did NOT find is as important as what we found"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.2), VL_RED)
tb(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.3), "H2.2 REJECTED: NOT Substrate-Specific", sz=15, bold=True, color=ACCENT_RED)
bullets(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(1.6), [
    "Mann-Whitney U comparing N-C diff between groups:",
    "  GroEL substrates vs E. coli bg: p = 0.058 (NS)",
    "  HSP60 substrates vs mito bg: p = 0.536 (NS)",
    "  HSP60 substrates vs matrix bg: p = NS",
    "Background proteins show the SAME N>C asymmetry",
], sz=12, sp=Pt(2))

add_rect(slide, Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.2), VL_RED)
tb(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.3), "H2.3 REJECTED: No GroEL Class Gradient", sz=15, bold=True, color=ACCENT_RED)
bullets(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(1.6), [
    "Kruskal-Wallis H across Class I/II/III:",
    "  Contact order diff: p = 0.77, eta^2 = -0.014",
    "  pLDDT diff: p = 0.92",
    "If asymmetry drove dependence, Class III (obligate)",
    "should show greatest effect. They DON'T.",
], sz=12, sp=Pt(2))

fig3 = os.path.join(FIG_DIR, "fig3_groel_class_comparison.png")
if os.path.exists(fig3): slide.shapes.add_picture(fig3, Inches(0.3), Inches(3.7), Inches(12.7), Inches(3.0))

tb(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
   "Biological interpretation: N>C asymmetry is a 'gravitational constant' of protein architecture — reflects co-translational folding physics (N-terminus synthesized first, folds first), NOT chaperonin biology",
   sz=12, bold=True, color=DARK_BLUE)

# ===================== SLIDE 23: MTS RESULTS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Results: Goal 3 — Mitochondrial Targeting Signal Architecture", "HSP60 substrates are matrix-enriched; MTS creates a 'landing pad' for HSP60"); footer(slide, sn)
fig4 = os.path.join(FIG_DIR, "fig4_mts_targeting.png")
if os.path.exists(fig4): slide.shapes.add_picture(fig4, Inches(0.2), Inches(1.2), Inches(12.9), Inches(3.3))

add_rect(slide, Inches(0.3), Inches(4.7), Inches(6.2), Inches(1.6), VL_GREEN)
tb(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(0.3), "H3.1 CONFIRMED: Matrix Enrichment", sz=14, bold=True, color=ACCENT_GREEN)
bullets(slide, Inches(0.5), Inches(5.15), Inches(5.8), Inches(1.0), [
    "HSP60 substrates: 46.6% matrix | Background: 28.5%",
    "OR = 3.29, 95% CI [2.46, 4.40], p = 1.6x10^-16",
    "21.1% non-canonical matrix import (no detectable MTS)",
], sz=12, sp=Pt(2))

add_rect(slide, Inches(6.8), Inches(4.7), Inches(6.2), Inches(1.6), VL_GREEN)
tb(slide, Inches(7.0), Inches(4.8), Inches(5.8), Inches(0.3), "H3.3 CONFIRMED: MTS = Pre-Domain Extension", sz=14, bold=True, color=ACCENT_GREEN)
bullets(slide, Inches(7.0), Inches(5.15), Inches(5.8), Inches(1.0), [
    "84.4% (368/436) transit peptides end BEFORE first domain",
    "Binomial test (H0: p=0.5): p = 3.4x10^-51",
    "Median gap: 18 residues | Mean: 37.4 | Range: 0-579",
], sz=12, sp=Pt(2))

add_rect(slide, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.7), VL_GREEN)
tb(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
   '"Landing Pad" Model: After mitochondrial import & MTS cleavage, protein emerges with ~18aa unstructured linker -> unfolded first domain = optimally positioned for HSP60 capture. Architecture is NOT coincidental — it facilitates chaperonin-assisted folding upon import.',
   sz=13, bold=True, color=ACCENT_GREEN)

# ===================== SLIDE 24: CROSS-SPECIES CONSERVATION =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Results: Cross-Species Conservation — 2 Billion Years", "N-domain complexity conserved (r=0.84); 79.7% same CATH fold"); footer(slide, sn)
fig5 = os.path.join(FIG_DIR, "fig5_orthology.png")
if os.path.exists(fig5): slide.shapes.add_picture(fig5, Inches(0.2), Inches(1.2), Inches(12.9), Inches(3.5))
bullets(slide, Inches(0.5), Inches(4.9), Inches(12), Inches(2.3), [
    "69 cross-species homolog pairs: 33 both methods (highest confidence), 7 RBH-only, 29 orthogroup-only",
    "N-domain contact order conservation: Spearman r = 0.84, p = 5.3x10^-13 (n=45 paired values)",
    "    -> Topological complexity of N-terminal domain STRONGLY conserved across 2 billion years",
    "Fold type conservation: 55/69 pairs (79.7%) share the SAME top CATH superfamily",
    "RBH class observation: Only 8/84 (9.5%) Class III obligate have HSP60 orthologs, vs 18.4% Class I, 19.0% Class II",
    "    -> Most GroEL-dependent proteins are the most evolutionarily divergent from human counterparts",
    "Interpretation: selective pressure maintaining N-terminal folding complexity is ancient and strong",
], sz=13, sp=Pt(4))

# ===================== SLIDE 25: SUMMARY FIGURE =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Results Summary — All 56 Tests", "25 significant after hierarchical BH correction (FDR < 0.05)"); footer(slide, sn)
fig6 = os.path.join(FIG_DIR, "fig6_summary.png")
if os.path.exists(fig6): slide.shapes.add_picture(fig6, Inches(0.2), Inches(1.2), Inches(12.9), Inches(3.5))
stat = [
    ["Family", "Total Tests", "Significant", "Discovery Rate", "Key Result"],
    ["Domain Architecture", "24", "9", "37.5%", "TIM barrel OR=8.4; 1.10.10.10 OR=50.9"],
    ["N-vs-C Asymmetry", "30", "14", "46.7%", "N>C universal (r=0.41-0.48); substrate vs bg = NS"],
    ["MTS Targeting", "2", "2", "100%", "Matrix OR=3.29; pre-domain 84.4%"],
    ["TOTAL", "56", "25", "44.6%", "Hierarchical BH, within + across family correction"],
]
table(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(1.8), stat)
tb(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
   "Phase 1 pilot (281 tests, 22 significant) -> Phase 2 full-scale (56 tests, 25 significant): consolidated to pre-registered hypotheses, gained power",
   sz=11, color=GRAY, align=PP_ALIGN.CENTER)

# ===================== SLIDE 26: BIOLOGICAL SYNTHESIS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Biological Synthesis — Three Goals Converge", "What makes a chaperonin substrate? How does the mito system work? What's conserved?"); footer(slide, sn)
for i, (title, color, bg, items, lp) in enumerate([
    ("What Makes a Chaperonin\nSubstrate?", MED_BLUE, VL_BLUE, [
        "It's the FOLD, not general difficulty",
        "TIM barrels (OR=8.4): can't fold co-translationally",
        "8-fold symmetry needs all elements present",
        "NOT about domain count (OR=1.13, NS)",
        "NOT about N-vs-C asymmetry (substrate vs bg: NS)",
        "NOT about obligate dependence (KW p=0.77)",
    ], Inches(0.3)),
    ("How Does the Mitochondrial\nSystem Work?", ACCENT_ORANGE, VL_ORANGE, [
        "HSP60 substrates = matrix residents (3.3x OR)",
        "MTS creates a 'landing pad'",
        "84.4% pre-domain extension (p=3.4x10^-51)",
        "~18 residue gap after MTS cleavage",
        "First domain emerges unfolded in matrix",
        "Optimal positioning for HSP60 capture",
    ], Inches(4.5)),
    ("What is Conserved Across\n2 Billion Years?", ACCENT_GREEN, VL_GREEN, [
        "The fold itself: 79.7% same CATH SF",
        "N-domain complexity: r = 0.84 (p=5.3x10^-13)",
        "TIM barrels need GroEL in E. coli,",
        "    same folds need HSP60 in humans",
        "N>C asymmetry = universal property",
        "Translational physics conserved across all life",
    ], Inches(8.7)),
]):
    add_rect(slide, lp, Inches(1.3), Inches(4.0), Inches(2.8), bg)
    tb(slide, lp+Inches(0.2), Inches(1.35), Inches(3.6), Inches(0.4), title, sz=13, bold=True, color=color)
    bullets(slide, lp+Inches(0.2), Inches(1.85), Inches(3.6), Inches(2.1), items, sz=11, sp=Pt(2))

add_rect(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.8), DARK_BLUE)
tb(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.35), '"The End is the Beginning" — Explained', sz=18, bold=True, color=ACCENT_ORANGE)
bullets(slide, Inches(0.5), Inches(4.85), Inches(12), Inches(2.1), [
    "The C-terminus ('end' of translation) IS where the chaperonin's work begins — but NOT because C-regions are more complex",
    "The entire protein, synthesized N-to-C, has its most topologically complex domain at the N-terminus",
    "The N>C contact order asymmetry is a footprint of co-translational folding history, conserved across ALL multi-domain proteins",
    "Chaperonin substrate identity is determined by FOLD TOPOLOGY (TIM barrels, winged helix), not positional complexity",
    "The chaperonin captures the ENTIRE protein and helps the globally complex topology snap into place",
    "This is a 'gravitational constant' of protein architecture — always present, not caused by chaperonins",
], sz=13, color=WHITE, sp=Pt(3))

# ===================== SLIDE 27: LIMITATIONS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Limitations and Caveats"); footer(slide, sn)
for i, (title, items, lp) in enumerate([
    ("Methodological", [
        "pLDDT = confidence, not stability (mitigated by CO + FoldX)",
        "Co-IP captures interactions, not function (mitigated by SILAC>5)",
        "AlphaFold structures are predictions (mean pLDDT 85-93)",
        "Contact order imperfect for multi-state folders (r~-0.75 for two-state only)",
        "MitoCarta annotations evolve (70 reclassifications v2->v3)",
    ], Inches(0.3)),
    ("Statistical", [
        "56 tests: some false positive risk despite BH correction",
        "Only 69 homolog pairs — limited power for cross-species subtle effects",
        "10 kDa size-matching bins are approximate at distribution extremes",
        "No Class III equivalents for HSP60 (no dependence classification)",
        "Seed=42 for size-matched sampling (sensitivity analysis needed)",
    ], Inches(4.5)),
    ("Biological & Pending", [
        "~30-40% of matrix proteins lack detectable MTS",
        "No TargetP 2.0 / SignalP6 predictions (DTU license required)",
        "Alternative import pathways not characterized",
        "FoldX DeltaG: ~42% complete (est. April 1-2, 2026)",
        "    -> Will add thermodynamic stability dimension",
        "    -> Re-run Modules F, H, I after completion",
    ], Inches(8.7)),
]):
    add_rect(slide, lp, Inches(1.3), Inches(4.0), Inches(5.5), VL_RED if i < 2 else VL_ORANGE)
    tb(slide, lp+Inches(0.2), Inches(1.4), Inches(3.6), Inches(0.3), title, sz=14, bold=True, color=ACCENT_RED if i < 2 else ACCENT_ORANGE)
    bullets(slide, lp+Inches(0.2), Inches(1.8), Inches(3.6), Inches(4.8), items, sz=11, sp=Pt(3))

# ===================== SLIDE 28: FUTURE WORK =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Future Directions"); footer(slide, sn)
bullets(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Immediate (Session 7, ~April 2026):",
    "  1. Verify FoldX completion (25,007 proteins)",
    "  2. Submit 07_foldx_collect.sh to merge results",
    "  3. Re-run Module F with DeltaG integration",
    "  4. Re-run Module H: test N-domain vs C-region DeltaG",
    "  5. Re-run Module I: add DeltaG violin plots",
    "  6. Transfer final results to Mac",
    "  7. Manuscript preparation",
    "",
    "Medium-term:",
    "  8. TargetP 2.0 MTS predictions (DTU license)",
    "  9. IUPred2A intrinsic disorder in pre-domain tails",
    "  10. SignalP6 for signal peptide refinement",
    "  11. Sensitivity analysis with different SILAC cutoffs",
], sz=13, sp=Pt(3))
bullets(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Longer-term:",
    "  12. Extend to Group II chaperonins (TRiC/CCT)",
    "       Eukaryotic cytoplasmic chaperonin (8-subunit ring)",
    "       Would test if fold topology preference is universal",
    "  13. MD validation of FoldX DeltaG predictions",
    "       Select top 10 most interesting proteins",
    "       100 ns all-atom simulations per protein",
    "  14. Deep learning substrate prediction model",
    "       Features: fold type, CO, DeltaG, domain count",
    "       Train on GroEL, validate on HSP60",
    "  15. Experimental validation of 'landing pad' model",
    "       Truncation experiments: remove pre-domain linker",
    "       Test if HSP60 engagement efficiency decreases",
    "  16. Single-molecule folding studies",
    "       Optical tweezers on TIM barrel substrates",
    "       Test co-translational vs post-translational folding",
], sz=13, sp=Pt(3))

# ===================== SLIDE 29: SOFTWARE & REPRODUCIBILITY =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Reproducibility & Data Inventory"); footer(slide, sn)
bullets(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(3.0), [
    "All code organized in 3 locations:",
    "  scripts/ — 4 standalone cleaning scripts",
    "  workflow/scripts/ — 16 Phase 1 module scripts",
    "  workflow/phase2/ — Phase 2 pipeline:",
    "      Snakefile + config.yaml + 7 Python scripts",
    "      slurm_jobs/ — 19 SLURM scripts",
    "",
    "Configuration: workflow/phase2/config.yaml",
    "  All paths, parameters, SLURM resources in one file",
    "  HPC deployment: docs/HPC_DEPLOYMENT_GUIDE.md",
], sz=12, sp=Pt(2))
bullets(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(3.0), [
    "Data scale:",
    "  Phase 1: 1,390 proteins, ~466 MB structures",
    "  Phase 2: 25,007 proteins, ~22 GB structures",
    "  Total results: ~186 MB (transferred to Mac)",
    "  6 publication figures (PDF + PNG, 300 DPI)",
    "",
    "Raw data: all downloadable from public databases",
    "  UniProt, AlphaFold DB, MitoCarta 3.0",
    "  No proprietary data dependencies",
    "  Random seed = 42 for all stochastic operations",
], sz=12, sp=Pt(2))

inv = [
    ["Category", "Phase 1 (Pilot)", "Phase 2 (Full-Scale)"],
    ["Proteins analyzed", "1,390", "25,007"],
    ["Domain assignments", "1,387 (99.8%)", "25,258 unified records"],
    ["Foldseek clusters", "1,155 (24 shared)", "16,193 (27,063 proteins)"],
    ["N-vs-C paired", "~1,200", "2,648 pairs"],
    ["Contact order records", "~4,000", "11,824"],
    ["Statistical tests", "281 (22 significant)", "56 (25 significant)"],
    ["Publication figures", "6 (PDF + PNG)", "6 (polished, PDF + PNG)"],
    ["FoldX stability", "-", "~42% complete (25,007 target)"],
]
table(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(2.8), inv)

# ===================== SLIDE 30: SESSION HISTORY =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Development History — 6 Sessions"); footer(slide, sn)
hist = [
    ["Session", "Date", "Key Accomplishments", "Bugs Fixed / Lessons Learned"],
    ["1", "2026-03-10", "Project planning, GroEL standardization (149 demerged),\nHSP60 SILAC filtering (325->266), MitoCarta parsing", "NDIC = highest enrichment (not missing data)"],
    ["2", "2026-03-11", "FASTA extraction, RBH (40 pairs), AlphaFold download\n(1,382 CIFs), DSSP, CATH domains, Chainsaw, OrthoFinder", "API rate limiting (1 req/s for InterPro)"],
    ["3", "2026-03-12", "Modules F (N-vs-C), G (MTS), H (statistics),\nI (figures). Phase 1 COMPLETE.", "pLDDT != stability (critical insight)"],
    ["4", "2026-03-14", "Phase 2 pipeline, HPC deployment, Chainsaw full-scale.\n19 SLURM scripts, Snakemake workflow.", "bioconda stride != protein STRIDE (hours lost)\nColumn name mismatches (6 wrong lookups)\nHuman CIFs still .gz compressed"],
    ["5", "2026-03-17", "Foldseek full-scale (16K clusters), Module E verified,\nF->H->I chain completed. Results transferred to Mac.", "LD_LIBRARY_PATH fix for gcc-8.4/numpy\npython3 -u for SLURM real-time logs"],
    ["6", "2026-03-22", "FoldX 5.1 installed + tested. Array job submitted\n(501 tasks). 2 timeout fixes. Collaborator deliverables.", "FoldX 5.1: component energies null (OK, total works)\nTimeout chunks 125,139 -> extended to 6h"],
    ["7", "(pending)", "FoldX completion, Module F/H/I re-run with DeltaG,\nManuscript preparation.", ""],
]
table(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(5.8), hist,
      cw=[Inches(0.7), Inches(1.2), Inches(5.5), Inches(5.5)])

# ===================== SLIDE 31: KEY BUG FIXES =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Key Bug Fixes & Lessons Learned", "Critical issues discovered and resolved during development"); footer(slide, sn)

add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.0), Inches(2.8), VL_RED)
tb(slide, Inches(0.5), Inches(1.4), Inches(3.6), Inches(0.3), "Bug 1: Wrong STRIDE Binary", sz=14, bold=True, color=ACCENT_RED)
bullets(slide, Inches(0.5), Inches(1.8), Inches(3.6), Inches(2.2), [
    "conda install stride -> StriDe v0.0.1",
    "    (genomic structural variant caller!)",
    "Protein STRIDE = Frishman & Argos (1995)",
    "    Must compile from heiniglab/stride",
    "Impact: ALL Chainsaw predictions failed",
    "    -> ndom=0 for 4,457 proteins",
    "    Zero output, no error message",
    "Fix: Compiled correct STRIDE from source",
], sz=11, sp=Pt(1))

add_rect(slide, Inches(4.5), Inches(1.3), Inches(4.0), Inches(2.8), VL_RED)
tb(slide, Inches(4.7), Inches(1.4), Inches(3.6), Inches(0.3), "Bug 2: Human CIFs Compressed", sz=14, bold=True, color=ACCENT_RED)
bullets(slide, Inches(4.7), Inches(1.8), Inches(3.6), Inches(2.2), [
    "EBI bulk download: human CIFs as .cif.gz",
    "    E. coli was already uncompressed",
    "Impact: only 4,457 structures found",
    "    (should be 25,007)",
    "Fix: 05a_decompress_human.sh",
    "    gunzip 23,672 human CIF files",
    "Lesson: always verify file format",
    "    after bulk downloads",
], sz=11, sp=Pt(1))

add_rect(slide, Inches(8.7), Inches(1.3), Inches(4.3), Inches(2.8), VL_RED)
tb(slide, Inches(8.9), Inches(1.4), Inches(3.9), Inches(0.3), "Bug 3: Column Name Mismatches", sz=14, bold=True, color=ACCENT_RED)
bullets(slide, Inches(8.9), Inches(1.8), Inches(3.9), Inches(2.2), [
    "6 wrong column lookups in Module E:",
    "  CATH uses 'uniprot_accession'",
    "  GroEL uses 'current_accession'",
    "  HSP60 uses 'uniprot_id'",
    "  Matrix/Mito also use 'uniprot_id'",
    "Impact: empty accession columns,",
    "    zero substrate counts",
    "Fix: Rewrote with correct names +",
    "    fallback detection logic",
], sz=11, sp=Pt(1))

# Additional lessons
add_rect(slide, Inches(0.3), Inches(4.3), Inches(12.7), Inches(2.9), RGBColor(0xF5,0xF5,0xF5))
tb(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(0.3), "HPC-Specific Lessons", sz=14, bold=True, color=DARK_BLUE)
bullets(slide, Inches(0.5), Inches(4.8), Inches(6.0), Inches(2.2), [
    "LD_LIBRARY_PATH fix (CRITICAL):",
    "  HPC gcc-8.4 has old libstdc++.so",
    "  Conda numpy/scipy need newer version",
    "  Every SLURM script MUST include:",
    "  export LD_LIBRARY_PATH=$CONDA_PREFIX/lib:$LD_LIBRARY_PATH",
    "  Without this: silent ImportError on Python import",
], sz=11, sp=Pt(2))
bullets(slide, Inches(6.8), Inches(4.8), Inches(6.0), Inches(2.2), [
    "Python unbuffered output:",
    "  SLURM captures stdout at job END, not real-time",
    "  All scripts use: python3 -u (unbuffered)",
    "  Enables: tail -f slurm-*.out for monitoring",
    "",
    "FoldX timeout: 4h -> 6h for large proteins",
    "  Chunks 125, 139 had large proteins (>1000 aa)",
    "  Resubmitted as job 94439 with 6h wall time",
], sz=11, sp=Pt(2))

# ===================== SLIDE 32: COMPLETE FILE INVENTORY =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Complete File Inventory — Key Outputs", "Every major input, result, and script file"); footer(slide, sn)

inv2 = [
    ["Category", "File", "Records", "Description"],
    ["Input", "ecoli_k12_proteome.fasta", "4,403", "E. coli K-12 reference proteome"],
    ["Input", "human_proteome.fasta", "20,416", "Human canonical proteome"],
    ["Input", "groel_substrates_standardized.tsv", "252", "GroEL substrates (16 columns)"],
    ["Input", "hsp60_tier1_substrates.tsv", "266", "HSP60 Tier-1 (28 columns)"],
    ["Input", "human_mito_proteome.tsv", "1,136", "MitoCarta 3.0 full mito"],
    ["Input", "human_matrix_proteome.tsv", "525", "Matrix-only subset"],
    ["Input", "groel_hsp60_homologs.tsv", "69", "Dataset 6 homolog pairs"],
    ["Result", "structure_index.tsv", "1,382", "AlphaFold metadata + pLDDT"],
    ["Result", "dssp_summary.tsv", "1,382", "Per-protein SS composition"],
    ["Result", "cath_domain_assignments.tsv", "2,141", "Gene3D domain records"],
    ["Result", "ml_domain_assignments.tsv", "1,387", "Unified CATH+Chainsaw (pilot)"],
    ["Phase 2", "unified_domain_assignments_full.tsv", "25,258", "All domain assignments"],
    ["Phase 2", "n_vs_c_paired_full.tsv", "2,648", "Paired N-vs-C comparisons"],
    ["Phase 2", "contact_order_full.tsv", "11,824", "Per-region RCO"],
    ["Phase 2", "corrected_pvalues_full.tsv", "56", "All Phase 2 statistical tests"],
    ["Phase 2", "foldseek_clusters_full.tsv", "16,193", "Structural clusters"],
    ["Figures", "fig[1-6]_*.{pdf,png}", "12 files", "6 publication figures, 300 DPI"],
]
table(slide, Inches(0.1), Inches(1.2), Inches(13.1), Inches(5.8), inv2,
      cw=[Inches(0.9), Inches(3.5), Inches(1.0), Inches(7.7)])

# ===================== SLIDE 33: CONCLUSIONS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, DARK_BLUE)
tb(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5), "Conclusions", sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
bullets(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(6.0), [
    "1. Chaperonin substrates are enriched for SPECIFIC complex fold topologies",
    "      GroEL: TIM barrels (OR=8.4, p=1.9x10^-8), Winged helix (OR=50.9, p=3.5x10^-9)",
    "      HSP60: Rossmann-like (OR=5.4), beta sandwiches (OR=3.6)",
    "      NOT enriched for multi-domain proteins (OR=1.13, NS)",
    "",
    "2. N-terminal domains universally have higher contact order than C-terminal regions",
    "      Effect is UNIVERSAL — present in ALL multi-domain proteins (strongest in mito bg: r=0.48)",
    "      NOT substrate-specific (Mann-Whitney GroEL vs bg: p=0.058; HSP60 vs bg: p=0.536)",
    "      NO GroEL class gradient (Kruskal-Wallis p=0.77) — Class III = Class I asymmetry",
    "",
    "3. Mitochondrial targeting signals create a 'landing pad' for post-import folding",
    "      84.4% pre-domain (p=3.4x10^-51), median gap=18 residues",
    "      HSP60 substrates 3.3x enriched for matrix (p=1.6x10^-16)",
    "",
    "4. Chaperonin substrate properties are CONSERVED across 2 billion years",
    "      N-domain contact order: r=0.84 (p=5.3x10^-13) | 79.7% same CATH superfamily",
    "",
    "5. The N>C asymmetry is a 'gravitational constant' of protein architecture",
    "      Reflects co-translational folding physics (N-terminus synthesized first), not chaperonin biology",
], sz=14, color=WHITE, sp=Pt(3))

# ===================== SLIDE 34: KEY NUMBERS =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
header(slide, "Key Numbers at a Glance"); footer(slide, sn)
nums = [
    ["Metric", "Value", "Context"],
    ["Total proteins analyzed", "25,007", "4,403 E. coli + 20,416 human"],
    ["GroEL substrates", "252 (38/126/84 by class)", "Kerner et al. 2005"],
    ["HSP60 Tier-1 substrates", "266", "SILAC > 5, MitoCarta confirmed"],
    ["Cross-species homologs", "69 pairs", "33 both + 7 RBH + 29 orthogroup"],
    ["Domain coverage", "99.8% (pilot), 93.6% (full)", "CATH + Chainsaw combined"],
    ["Foldseek clusters", "16,193 (75.6% singletons)", "3Di+AA structural alphabet"],
    ["Statistical tests", "56 total, 25 significant", "Hierarchical BH, FDR<0.05"],
    ["TIM barrel enrichment (GroEL)", "OR = 8.4, p = 1.9x10^-8", "Most robust domain finding"],
    ["Winged helix enrichment (GroEL)", "OR = 50.9, p = 3.5x10^-9", "Strongest individual enrichment"],
    ["N>C contact order (strongest)", "r = 0.48, p = 7.1x10^-18", "Mito background (UNIVERSAL)"],
    ["Substrate vs background N>C", "p = 0.536 (NS)", "NOT substrate-specific"],
    ["GroEL class gradient", "p = 0.77 (NS)", "No Class III > Class I effect"],
    ["Matrix enrichment (HSP60)", "OR = 3.29, p = 1.6x10^-16", "HSP60 substrates -> matrix"],
    ["MTS pre-domain", "84.4%, p = 3.4x10^-51", "Median gap = 18 residues"],
    ["N-domain RCO conservation", "r = 0.84, p = 5.3x10^-13", "Across 2 billion years"],
    ["FoldX progress", "~42% (10,775/25,007)", "Est. completion April 1-2"],
]
table(slide, Inches(0.2), Inches(1.2), Inches(12.9), Inches(5.8), nums,
      cw=[Inches(3.5), Inches(3.0), Inches(6.4)])

# ===================== SLIDE 35: THANK YOU =====================
sn += 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06), ACCENT_ORANGE)
tb(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.7), "Thank You", sz=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(0.5), "Antah Asti Prarambh — The End is the Beginning", sz=22, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.4), "Vishal Bharti", sz=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.4), "CSIR-Institute of Genomics and Integrative Biology (IGIB), New Delhi", sz=15, color=RGBColor(0xBB,0xCC,0xDD), align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(2.0),
   "Pipeline: 16 Python modules + 7 Phase 2 scripts + 19 SLURM jobs\n"
   "Data: 25,007 proteins | 7 datasets | 56 statistical tests | 25 significant\n"
   "Tools: MMseqs2, Foldseek, Chainsaw, STRIDE, FoldX, mkDSSP, AlphaFold DB\n"
   "Figures: 6 publication-quality (PDF+PNG, 300 DPI, colorblind-friendly)\n"
   "Status: Phase 2 complete; FoldX ~42% done; manuscript preparation next",
   sz=15, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)

# SAVE
prs.save(OUT_FILE)
print(f"Presentation saved to: {OUT_FILE}")
print(f"Total slides: {sn}")
assert sn == TOTAL_SLIDES, f"Slide count mismatch: {sn} != {TOTAL_SLIDES}"
